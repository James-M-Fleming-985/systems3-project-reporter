#!/usr/bin/env python3
"""ZnNi Line Report Generator - CLI Application."""
import argparse
import importlib.util
import json
import logging
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional

from src.utils import setup_logging, load_config, ensure_output_directory, generate_output_filename, validate_input_files

# Feature folder mappings
FEATURE_FOLDERS = {
    "data_reader": "FEATURE-003-001_Data_Reader_Parser",
    "risk_aggregator": "FEATURE-003-002_Risk_Aggregator",
    "gantt_chart": "FEATURE-003-003_Gantt_Chart_Generator",
    "milestone_tracker": "FEATURE-003-004_Milestone_Tracker",
    "change_logger": "FEATURE-003-005_Change_Management_Logger",
    "powerpoint_generator": "FEATURE-003-006_PowerPoint_Generator"
}


def load_feature_orchestrator(feature_folder_name: str):
    """Load feature orchestrator module dynamically."""
    path = Path(__file__).parent / feature_folder_name / "src" / "feature_integration.py"
    if not path.exists():
        raise FileNotFoundError(f"Feature integration module not found: {path}")
    
    spec = importlib.util.spec_from_file_location(f"{feature_folder_name}.integration", path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module




def is_success(response):
    """Check if response indicates success - handles both response formats."""
    if hasattr(response, 'status'):
        # Data Reader format: status enum
        return response.status.value == "success"
    elif hasattr(response, 'success'):
        # Risk Aggregator format: success bool
        return response.success
    return False





class ZnNiReportGenerator:
    """Main report generator orchestrating all features."""
    
    def __init__(self, config_path: str = "config/settings.yaml"):
        self.config = load_config(config_path)
        self.logger = setup_logging(self.config.get("data_reader", {}).get("log_level", "INFO"))
        self.output_dir = ensure_output_directory(self.config.get("output", {}).get("base_dir", "output"))
        self.orchestrators = {}
        self._initialize_features()
    
    def _initialize_features(self):
        """Initialize all feature orchestrators."""
        for feature_key, folder_name in FEATURE_FOLDERS.items():
            try:
                module = load_feature_orchestrator(folder_name)
                feature_config_dict = self.config.get(feature_key, {})
                
                if feature_config_dict and hasattr(module, 'FeatureConfig'):
                    feature_config = module.FeatureConfig(**feature_config_dict)
                    orchestrator = module.FeatureOrchestrator(feature_config)
                else:
                    orchestrator = module.FeatureOrchestrator()
                
                self.orchestrators[feature_key] = orchestrator
                self.logger.info(f"Initialized feature: {feature_key}")
            except Exception as e:
                self.logger.error(f"Failed to initialize feature {feature_key}: {e}")
                raise
    
    def process_data_files(self, input_files: list) -> Dict[str, Any]:
        """Process input data files using Data Reader Parser."""
        self.logger.info(f"Processing {len(input_files)} input files")
        data_reader = self.orchestrators["data_reader"]
        
        all_data = []
        for file_path in input_files:
            response = data_reader.read_and_validate(file_path)
            if is_success(response):
                all_data.append(response.data)
            else:
                self.logger.warning(f"Failed to read {file_path}: {response.errors}")
        
        return {"raw_data": all_data, "parsed_data": all_data, "file_count": len(all_data)}
    
    def aggregate_risks(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Aggregate risk data using Risk Aggregator."""
        self.logger.info("Aggregating risk data")
        risk_aggregator = self.orchestrators["risk_aggregator"]
        
        # Extract risks from raw data (embedded in YAML)
        all_risks = []
        for item in data.get("raw_data", []):
            if "risks" in item:
                all_risks.extend(item["risks"])
        
        # For now, just return the risks - Risk Aggregator expects file paths
        # In a real scenario, we would save risks to temp file or refactor the feature
        return {
            "risks": all_risks,
            "summary": {"total": len(all_risks)}
        }
    def track_milestones(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Track project milestones."""
        self.logger.info("Tracking milestones")
        milestone_tracker = self.orchestrators["milestone_tracker"]
        
        milestones = []
        for item in data.get("raw_data", []):
            milestone_data = item.get("milestones", [])
            if milestone_data:
                response = milestone_tracker.batch_track_milestones(milestone_data)
                if is_success(response):
                    milestones.extend(response.data.get("tracked", []))
        
        # Create milestone quadrant
        quadrant_response = milestone_tracker.create_milestone_quadrant(milestones)
        return {
            "milestones": milestones,
            "quadrant": quadrant_response.data if is_success(response) else {}
        }
    
    def generate_gantt_chart(self, milestone_data: Dict[str, Any]) -> Optional[str]:
        """Generate Gantt chart from milestones."""
        self.logger.info("Generating Gantt chart")
        gantt_generator = self.orchestrators["gantt_chart"]
        
        milestones = milestone_data.get("milestones", [])
        if not milestones:
            self.logger.warning("No milestones available for Gantt chart")
            return None
        
        output_path = self.output_dir / generate_output_filename(
            self.config["output"]["report_name_prefix"] + "_gantt",
            "png",
            self.config.get("output", {}).get("timestamp_format", "%Y%m%d_%H%M%S")
        )
        
        response = gantt_generator.generate_gantt_chart(milestones, str(output_path))
        if is_success(response):
            self.logger.info(f"Gantt chart saved to: {output_path}")
            return str(output_path)
        else:
            self.logger.error(f"Failed to generate Gantt chart: {response.errors}")
            return None
    
    def log_changes(self, changes: list) -> Dict[str, Any]:
        """Log change management records."""
        self.logger.info(f"Logging {len(changes)} changes")
        change_logger = self.orchestrators["change_logger"]
        
        if changes:
            response = change_logger.log_batch_changes(changes)
            if is_success(response):
                # Get log statistics
                stats_response = change_logger.get_log_statistics()
                return stats_response.data if is_success(response) else {}
        
        return {"logged_count": 0}
    
    def generate_powerpoint_report(self, report_data: Dict[str, Any]) -> str:
        """Generate comprehensive PowerPoint report with detailed project data."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        self.logger.info("Generating PowerPoint report")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "ZnNi Line Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        subtitle.text += f"Projects: {report_data.get('file_count', 0)}"
        
        bullet_slide_layout = prs.slide_layouts[1]
        
        # Get project data
        project_data = report_data.get('project_data', [])
        
        # Process each project
        for proj in project_data:
            project_name = proj.get('project_name', 'Unknown Project')
            
            # Project Overview Slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            title_shape.text = f"Project: {project_name}"
            tf = body_shape.text_frame
            tf.text = f"Code: {proj.get('project_code', 'N/A')}"
            tf.add_paragraph().text = f"Status: {proj.get('status', 'N/A')}"
            pct = proj.get('completion_percentage', 0)
            tf.add_paragraph().text = f"Progress: {pct}%"
            end_date = proj.get('target_completion', 'N/A')
            tf.add_paragraph().text = f"Target: {end_date}"
            
            # Milestones Slide
            milestones = proj.get('milestones', [])
            if milestones:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                title_shape.text = f"{project_name} - Milestones"
                tf = body_shape.text_frame
                tf.text = f"Total Milestones: {len(milestones)}"
                
                # Group by status
                completed = [m for m in milestones if m.get('status') == 'COMPLETED']
                in_progress = [m for m in milestones if m.get('status') == 'IN_PROGRESS']
                not_started = [m for m in milestones if m.get('status') == 'NOT_STARTED']
                
                tf.add_paragraph().text = f"Completed: {len(completed)}"
                for m in completed[:3]:  # Show first 3
                    p = tf.add_paragraph()
                    mdate = m.get('target_date', 'N/A')
                    p.text = f"  • {m.get('name', 'N/A')} ({mdate})"
                    p.level = 1
                
                tf.add_paragraph().text = f"In Progress: {len(in_progress)}"
                for m in in_progress[:3]:
                    p = tf.add_paragraph()
                    mdate = m.get('target_date', 'N/A')
                    p.text = f"  • {m.get('name', 'N/A')} (Due: {mdate})"
                    p.level = 1
                
                tf.add_paragraph().text = f"Not Started: {len(not_started)}"
            
            # Risks Slide
            risks = proj.get('risks', [])
            if risks:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                title_shape.text = f"{project_name} - Risks"
                tf = body_shape.text_frame
                tf.text = f"Total Risks: {len(risks)}"
                
                # Group by severity
                high_risks = [r for r in risks if r.get('severity') == 'HIGH']
                medium_risks = [r for r in risks if r.get('severity') == 'MEDIUM']
                low_risks = [r for r in risks if r.get('severity') == 'LOW']
                
                if high_risks:
                    tf.add_paragraph().text = f"HIGH Severity: {len(high_risks)}"
                    for r in high_risks:
                        p = tf.add_paragraph()
                        p.text = f"  • {r.get('description', 'N/A')}"
                        p.level = 1
                
                if medium_risks:
                    tf.add_paragraph().text = f"MEDIUM Severity: {len(medium_risks)}"
                    for r in medium_risks:
                        p = tf.add_paragraph()
                        p.text = f"  • {r.get('description', 'N/A')}"
                        p.level = 1
                
                if low_risks:
                    tf.add_paragraph().text = f"LOW Severity: {len(low_risks)}"
            
            # Changes Slide
            changes = proj.get('changes', [])
            if changes:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                title_shape.text = f"{project_name} - Change Management"
                tf = body_shape.text_frame
                tf.text = f"Total Changes: {len(changes)}"
                
                for change in changes:
                    p = tf.add_paragraph()
                    p.text = f"• {change.get('change_id', 'N/A')}: {change.get('reason', 'N/A')}"
                    p2 = tf.add_paragraph()
                    p2.text = f"  Impact: {change.get('impact', 'N/A')}"
                    p2.level = 1
        
        # Overall Summary Slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        title_shape.text = "Overall Summary"
        tf = body_shape.text_frame
        
        total_milestones = sum(len(p.get('milestones', [])) for p in project_data)
        total_risks = sum(len(p.get('risks', [])) for p in project_data)
        total_changes = sum(len(p.get('changes', [])) for p in project_data)
        
        tf.text = f"Total Projects: {len(project_data)}"
        tf.add_paragraph().text = f"Total Milestones: {total_milestones}"
        tf.add_paragraph().text = f"Total Risks: {total_risks}"
        tf.add_paragraph().text = f"Total Changes: {total_changes}"
        
        # Save presentation
        output_path = Path("output") / "ZnNi_Report.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        self.logger.info(f"PowerPoint report saved to: {output_path}")
        return str(output_path)
    def generate_report(self, input_files: list, changes: list = None) -> Dict[str, Any]:
        """Generate complete ZnNi line report."""
        self.logger.info("Starting ZnNi Line Report generation")
        report_results = {}
        
        try:
            # Step 1: Process input data files
            data = self.process_data_files(input_files)
            report_results["file_count"] = data["file_count"]
            report_results["project_data"] = data.get("parsed_data", [])
            
            # Step 2: Aggregate risks
            risk_data = self.aggregate_risks(data)
            report_results["risk_summary"] = risk_data["summary"]
            
            # Step 3: Track milestones
            milestone_data = self.track_milestones(data)
            report_results["milestone_data"] = milestone_data
            
            # Step 4: Generate Gantt chart
            gantt_path = self.generate_gantt_chart(milestone_data)
            if gantt_path:
                report_results["gantt_chart"] = gantt_path
            
            # Step 5: Log changes if provided
            if changes:
                change_stats = self.log_changes(changes)
                report_results["change_stats"] = change_stats
            
            # Step 6: Generate PowerPoint report
            ppt_path = self.generate_powerpoint_report(report_results)
            report_results["powerpoint_report"] = ppt_path
            
            # Save JSON summary
            summary_path = self.output_dir / generate_output_filename(
                self.config["output"]["report_name_prefix"] + "_summary",
                "json",
                self.config.get("output", {}).get("timestamp_format", "%Y%m%d_%H%M%S")
            )
            with open(summary_path, 'w') as f:
                json.dump(report_results, f, indent=2, default=str)
            report_results["summary_file"] = str(summary_path)
            
            self.logger.info("Report generation completed successfully")
            return report_results
            
        except Exception as e:
            self.logger.error(f"Report generation failed: {e}")
            raise


def main():
    """CLI entry point."""
    parser = argparse.ArgumentParser(
        description="ZnNi Line Report Generator - Generate comprehensive reports from ZnNi line data"
    )
    parser.add_argument(
        "input_files",
        nargs="+",
        help="Input data files to process"
    )
    parser.add_argument(
        "--config",
        default="config/settings.yaml",
        help="Configuration file path (default: config/settings.yaml)"
    )
    parser.add_argument(
        "--changes",
        help="JSON file containing change management records"
    )
    parser.add_argument(
        "--output-dir",
        help="Override output directory from config"
    )
    parser.add_argument(
        "--log-level",
        choices=["DEBUG", "INFO", "WARNING", "ERROR"],
        help="Override log level from config"
    )
    
    args = parser.parse_args()
    
    try:
        # Validate input files
        input_files = validate_input_files(args.input_files)
        
        # Load change records if provided
        changes = []
        if args.changes:
            with open(args.changes, 'r') as f:
                changes = json.load(f)
        
        # Initialize generator
        generator = ZnNiReportGenerator(args.config)
        
        # Override settings if provided
        if args.output_dir:
            generator.output_dir = ensure_output_directory(args.output_dir)
        if args.log_level:
            generator.logger.setLevel(getattr(logging, args.log_level))
        
        # Generate report
        results = generator.generate_report(input_files, changes)
        
        # Print summary
        print("\n" + "=" * 60)
        print("ZnNi Line Report Generation Complete")
        print("=" * 60)
        print(f"Files Processed: {results.get('file_count', 0)}")
        print(f"PowerPoint Report: {results.get('powerpoint_report', 'N/A')}")
        print(f"Gantt Chart: {results.get('gantt_chart', 'N/A')}")
        print(f"Summary File: {results.get('summary_file', 'N/A')}")
        print("=" * 60)
        
    except Exception as e:
        print(f"\nError: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()