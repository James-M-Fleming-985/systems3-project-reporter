"""
End-to-end tests for PowerPoint Generator feature (FEATURE-003-006)
"""

import pytest
import os
import time
import tempfile
import shutil
from pathlib import Path
from datetime import datetime
from unittest.mock import patch, MagicMock
import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Assuming these are the main modules for the feature
from powerpoint_generator import PowerPointGenerator
from presentation_service import PresentationService
from template_manager import TemplateManager
from content_processor import ContentProcessor


class TestPowerPointGeneratorE2E:
    """End-to-end tests for PowerPoint Generator feature"""
    
    @pytest.fixture(autouse=True)
    def setup_and_teardown(self):
        """Setup test environment and cleanup after tests"""
        # Setup
        self.temp_dir = tempfile.mkdtemp()
        self.test_output_dir = Path(self.temp_dir) / "output"
        self.test_output_dir.mkdir(exist_ok=True)
        
        self.test_templates_dir = Path(self.temp_dir) / "templates"
        self.test_templates_dir.mkdir(exist_ok=True)
        
        # Create a sample template
        self._create_sample_template()
        
        yield
        
        # Teardown
        shutil.rmtree(self.temp_dir, ignore_errors=True)
    
    def _create_sample_template(self):
        """Create a sample PowerPoint template for testing"""
        prs = Presentation()
        
        # Title slide layout
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "{{title}}"
        subtitle.text = "{{subtitle}}"
        
        # Content slide layout
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "{{content_title}}"
        content.text = "{{content_body}}"
        
        template_path = self.test_templates_dir / "corporate_template.pptx"
        prs.save(str(template_path))
        return template_path
    
    @pytest.fixture
    def generator(self):
        """Create PowerPoint generator instance"""
        return PowerPointGenerator(
            template_dir=str(self.test_templates_dir),
            output_dir=str(self.test_output_dir)
        )
    
    @pytest.fixture
    def sample_presentation_data(self):
        """Sample data for creating a presentation"""
        return {
            "title": "Q4 2023 Financial Results",
            "subtitle": "Annual Performance Review",
            "author": "John Doe",
            "company": "Tech Corp Inc.",
            "template": "corporate_template.pptx",
            "slides": [
                {
                    "type": "title",
                    "title": "Q4 2023 Financial Results",
                    "subtitle": "Annual Performance Review"
                },
                {
                    "type": "content",
                    "title": "Executive Summary",
                    "content": [
                        "• Revenue increased by 25% YoY",
                        "• Operating margin improved to 18%",
                        "• Strong growth in cloud services",
                        "• Successful product launches in Q4"
                    ]
                },
                {
                    "type": "chart",
                    "title": "Revenue Growth",
                    "chart_type": "column",
                    "data": {
                        "categories": ["Q1", "Q2", "Q3", "Q4"],
                        "series": [
                            {"name": "2022", "values": [100, 110, 115, 120]},
                            {"name": "2023", "values": [120, 135, 145, 150]}
                        ]
                    }
                },
                {
                    "type": "table",
                    "title": "Key Metrics",
                    "headers": ["Metric", "2022", "2023", "Growth"],
                    "rows": [
                        ["Revenue ($M)", "445", "556", "+25%"],
                        ["Operating Margin", "15%", "18%", "+3pp"],
                        ["Customers", "1,200", "1,850", "+54%"]
                    ]
                },
                {
                    "type": "image",
                    "title": "Market Presence",
                    "image_url": "https://example.com/market_map.png",
                    "caption": "Global market presence as of Q4 2023"
                }
            ]
        }
    
    @pytest.mark.e2e
    def test_complete_presentation_generation_workflow(self, generator, sample_presentation_data):
        """
        Test complete end-to-end workflow for generating a PowerPoint presentation
        
        Scenario:
        1. User provides presentation data with multiple slide types
        2. System validates the input data
        3. System loads the specified template
        4. System generates all slides with proper formatting
        5. System saves the presentation to the output directory
        6. System returns the file path and metadata
        """
        # Act
        result = generator.create_presentation(sample_presentation_data)
        
        # Assert - Verify result structure
        assert result is not None
        assert "file_path" in result
        assert "file_name" in result
        assert "file_size" in result
        assert "slide_count" in result
        assert "generation_time" in result
        assert "status" in result
        
        # Verify file was created
        output_file = Path(result["file_path"])
        assert output_file.exists()
        assert output_file.suffix == ".pptx"
        
        # Verify file size is reasonable (not empty)
        assert result["file_size"] > 1000  # At least 1KB
        
        # Verify presentation content
        prs = Presentation(str(output_file))
        assert len(prs.slides) == len(sample_presentation_data["slides"])
        
        # Verify first slide (title slide)
        title_slide = prs.slides[0]
        assert sample_presentation_data["title"] in title_slide.shapes.title.text
        
        # Verify slide types were processed correctly
        for idx, slide_data in enumerate(sample_presentation_data["slides"][1:], 1):
            slide = prs.slides[idx]
            if slide_data["type"] == "content":
                # Verify content slide has title and body
                assert slide_data["title"] in slide.shapes.title.text
                # Check if content items are in the slide
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text
                for content_item in slide_data["content"]:
                    assert any(item in slide_text for item in content_item.split("•") if item.strip())
        
        # Verify metadata
        assert result["slide_count"] == len(sample_presentation_data["slides"])
        assert result["status"] == "completed"
        assert result["generation_time"] > 0
    
    @pytest.mark.e2e
    @patch('requests.get')
    def test_presentation_with_external_resources(self, mock_get, generator):
        """
        Test end-to-end workflow with external resources (images, data)
        
        Scenario:
        1. User creates presentation with external image URLs
        2. System downloads and embeds external images
        3. System handles network delays and retries
        4. System generates presentation with all resources
        5. System handles missing resources gracefully
        """
        # Mock external image request
        mock_image_response = MagicMock()
        mock_image_response.status_code = 200
        mock_image_response.content = b"fake_image_data"
        mock_image_response.headers = {"content-type": "image/png"}
        mock_get.return_value = mock_image_response
        
        presentation_data = {
            "title": "Product Launch 2024",
            "subtitle": "Innovation Showcase",
            "author": "Jane Smith",
            "template": "corporate_template.pptx",
            "slides": [
                {
                    "type": "title",
                    "title": "Product Launch 2024",
                    "subtitle": "Innovation Showcase"
                },
                {
                    "type": "image_gallery",
                    "title": "Product Screenshots",
                    "images": [
                        {
                            "url": "https://api.example.com/images/product1.png",
                            "caption": "Dashboard View"
                        },
                        {
                            "url": "https://api.example.com/images/product2.png",
                            "caption": "Analytics Module"
                        },
                        {
                            "url": "https://api.example.com/images/product3.png",
                            "caption": "Mobile Interface"
                        }
                    ]
                },
                {
                    "type": "data_visualization",
                    "title": "Market Analysis",
                    "data_source": "https://api.example.com/market-data",
                    "visualization_type": "pie_chart"
                }
            ],
            "branding": {
                "logo_url": "https://api.example.com/company-logo.png",
                "color_scheme": "blue",
                "footer_text": "Confidential - © 2024 Tech Corp Inc."
            }
        }
        
        # Act
        result = generator.create_presentation(presentation_data)
        
        # Assert
        assert result["status"] == "completed"
        assert mock_get.called
        assert mock_get.call_count >= 4  # 3 product images + 1 logo
        
        # Verify the presentation was created
        output_file = Path(result["file_path"])
        assert output_file.exists()
        
        # Verify presentation structure
        prs = Presentation(str(output_file))
        
        # Check if images were embedded (not checking actual image data in this test)
        image_slide = prs.slides[1]
        image_count = sum(1 for shape in image_slide.shapes 
                         if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        assert image_count > 0  # At least one image should be embedded
        
        # Verify branding elements
        for slide in prs.slides:
            footer_shapes = [shape for shape in slide.shapes 
                           if hasattr(shape, "text") and "Confidential" in shape.text]
            assert len(footer_shapes) > 0  # Footer should be on each slide
    
    @pytest.mark.e2e
    def test_error_handling_and_recovery(self, generator):
        """
        Test end-to-end error handling and recovery mechanisms
        
        Scenario:
        1. User provides invalid/incomplete data
        2. System validates and reports specific errors
        3. User provides data with unavailable template
        4. System falls back to default template
        5. User provides data causing processing errors
        6. System handles errors gracefully and provides meaningful feedback
        """
        # Test 1: Invalid input data (missing required fields)
        invalid_data = {
            "title": "Test Presentation",
            # Missing required fields: slides, template
        }
        
        with pytest.raises(ValueError) as exc_info:
            generator.create_presentation(invalid_data)
        assert "required field" in str(exc_info.value).lower()
        
        # Test 2: Invalid template
        data_with_invalid_template = {
            "title": "Test Presentation",
            "subtitle": "Test Subtitle",
            "template": "non_existent_template.pptx",
            "slides": [
                {
                    "type": "title",
                    "title": "Test",
                    "subtitle": "Subtitle"
                }
            ]
        }
        
        # Should fall back to default template or raise specific error
        result = generator.create_presentation(data_with_invalid_template)
        assert result["status"] in ["completed", "completed_with_warnings"]
        if "warnings" in result:
            assert any("template" in warning.lower() for warning in result["warnings"])
        
        # Test 3: Invalid slide type
        data_with_invalid_slide = {
            "title": "Test Presentation",
            "subtitle": "Test Subtitle",
            "template": "corporate_template.pptx",
            "slides": [
                {
                    "type": "title",
                    "title": "Test",
                    "subtitle": "Subtitle"
                },
                {
                    "type": "invalid_slide_type",
                    "title": "Invalid Slide"
                }
            ]
        }
        
        result = generator.create_presentation(data_with_invalid_slide)
        assert result["status"] in ["completed_with_warnings", "partial_success"]
        assert "warnings" in result or "errors" in result
        
        # Test 4: Very large presentation (stress test)
        large_presentation_data = {
            "title": "Large Presentation Test",
            "subtitle": "Stress Testing",
            "template": "corporate_template.pptx",
            "slides": [
                {
                    "type": "title",
                    "title": "Large Presentation Test",
                    "subtitle": "Stress Testing"
                }
            ] + [
                {
                    "type": "content",
                    "title": f"Slide {i}",
                    "content": [f"Content item {j}" for j in range(10)]
                }
                for i in range(50)  # 50 content slides
            ]
        }
        
        start_time = time.time()
        result = generator.create_presentation(large_presentation_data)
        end_time = time.time()
        
        assert result["status"] == "completed"
        assert result["slide_count"] == 51
        # Should complete within reasonable time (e.g., 30 seconds)
        assert (end_time - start_time) < 30
        
        # Test 5: Concurrent presentation generation
        import threading
        results = []
        errors = []
        
        def generate_presentation(data, index):
            try:
                result = generator.create_presentation(data)
                results.append((index, result))
            except Exception as e:
                errors.append((index, str(e)))
        
        # Create multiple presentations concurrently
        threads = []
        for i in range(3):
            data = {
                "title": f"Concurrent Test {i}",
                "subtitle": "Threading Test",
                "template": "corporate_template.pptx",
                "slides": [
                    {
                        "type": "title",
                        "title": f"Concurrent Test {i}",
                        "subtitle": "Threading Test"
                    }
                ]
            }
            thread = threading.Thread(target=generate_presentation, args=(data, i))
            threads.append(thread)
            thread.start()
        
        # Wait for all threads to complete
        for thread in threads:
            thread.join()
        
        # Verify all presentations were generated successfully
        assert len(errors) == 0
        assert len(results) == 3
        for index, result in results:
            assert result["status"] == "completed"
            assert Path(result["file_path"]).exists()


@pytest.mark.integration
class TestPowerPointGeneratorIntegration:
    """Integration tests for PowerPoint Generator with external services"""
    
    @pytest.fixture
    def service(self):
        """Create presentation service instance"""
        return PresentationService()
    
    @pytest.mark.e2e
    def test_template_management_workflow(self, service):
        """
        Test template upload, selection, and usage workflow
        
        Scenario:
        1. Admin uploads new presentation template
        2. System validates and stores template
        3. User lists available templates
        4. User selects template for presentation
        5. System generates presentation with selected template
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            # Create a custom template
            custom_template = Presentation()
            
            # Add custom layouts
            title_slide = custom_template.slides.add_slide(
                custom_template.slide_layouts[0]
            )
            title_slide.shapes.title.text = "{{custom_title}}"
            
            template_path = Path(temp_dir) / "custom_template.pptx"
            custom_template.save(str(template_path))
            
            # Upload template
            upload_result = service.upload_template(
                template_path=str(template_path),
                template_name="Marketing Template",
                category="marketing",
                description="Template for marketing presentations"
            )
            
            assert upload_result["status"] == "success"
            assert "template_id" in upload_result
            
            # List templates
            templates = service.list_templates(category="marketing")
            assert len(templates) > 0
            assert any(t["name"] == "Marketing Template" for t in templates)
            
            # Use template for presentation
            presentation_data = {
                "title": "Marketing Campaign 2024",
                "template_id": upload_result["template_id"],
                "slides": [
                    {
                        "type": "title",
                        "custom_title": "Marketing Campaign 2024"
                    }
                ]
            }
            
            result = service.generate_presentation(presentation_data)
            assert result["status"] == "completed"
            assert Path(result["file_path"]).exists()