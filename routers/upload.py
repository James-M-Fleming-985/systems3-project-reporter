"""
Upload Router - Handles XML file uploads and change management
"""
from fastapi import APIRouter, Request, UploadFile, File, Form, Depends
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.templating import Jinja2Templates
from pathlib import Path
import yaml
import os
import logging
from datetime import datetime
from typing import List

from services.xml_parser import MSProjectXMLParser
from services.change_detection import ChangeDetectionService
from services.subscription_service import SubscriptionService
from repositories.project_repository import ProjectRepository
from middleware.subscription import (
    get_user_or_create_anonymous, get_subscription_service, 
    enforce_upload_limits, SubscriptionError
)

router = APIRouter(tags=["upload"])
logger = logging.getLogger(__name__)

# Initialize services with persistent storage support
BASE_DIR = Path(__file__).resolve().parent.parent

# Use environment variables for persistent storage paths (Railway Volumes)
DATA_DIR = Path(os.getenv("DATA_STORAGE_PATH", str(BASE_DIR / "mock_data")))
UPLOAD_DIR = Path(os.getenv("UPLOAD_STORAGE_PATH", str(BASE_DIR / "uploads")))
TEMPLATES_DIR = BASE_DIR / "templates"
POWERPOINT_TEMPLATES_DIR = DATA_DIR / "powerpoint_templates"

# Ensure directories exist
DATA_DIR.mkdir(parents=True, exist_ok=True)
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
POWERPOINT_TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

project_repo = ProjectRepository(data_dir=DATA_DIR)
xml_parser = MSProjectXMLParser()
change_detector = ChangeDetectionService()
templates = Jinja2Templates(directory=str(TEMPLATES_DIR))


@router.get("/upload", response_class=HTMLResponse)
async def upload_page(request: Request):
    """Unified upload page for project XML and risk files"""
    from main import BUILD_VERSION
    
    # Get existing projects
    projects = project_repo.load_all_projects()
    
    # Get user from request state
    user = getattr(request.state, 'user', None) if hasattr(request, 'state') else None
    
    context = {
        "request": request,
        "projects": projects,
        "build_version": BUILD_VERSION,
        "user": user
    }
    
    return templates.TemplateResponse("upload_unified.html", context)



@router.post("/upload/xml")
async def upload_xml(
    file: UploadFile = File(...),
    is_baseline: str = Form("false"),
    clear_previous_changes: str = Form("false"),
    user=Depends(get_user_or_create_anonymous),
    sub_service: SubscriptionService = Depends(get_subscription_service)
):
    """
    Upload MS Project XML file and detect changes (with subscription limits)
    
    Args:
        clear_previous_changes: If "true", removes all previous changes before adding new ones
    """
    import logging
    logger = logging.getLogger(__name__)
    
    # Convert baseline flag to boolean
    is_baseline_upload = is_baseline.lower() == "true"
    
    try:
        logger.info(f"Starting XML upload for user {user.user_id}")
        logger.info(f"Baseline upload: {is_baseline_upload}")
        
        # Check file size and subscription limits
        file_size_mb = len(await file.read()) / (1024 * 1024)
        await file.seek(0)  # Reset file position
        
        logger.info(f"File size: {file_size_mb:.2f}MB")
        
        # Enforce upload limits
        try:
            enforce_upload_limits(file_size_mb, user, sub_service)
            logger.info("Upload limits check passed")
        except SubscriptionError as se:
            logger.warning(f"Subscription limit exceeded: {se.detail}")
            return JSONResponse({
                'success': False,
                'error': se.detail,
                'error_type': 'subscription_limit',
                'limit_info': se.limit_info if hasattr(se, 'limit_info') else None
            }, status_code=402)
        
        # Read and parse XML
        content = await file.read()
        xml_content = content.decode('utf-8')
        
        logger.info("XML file read successfully")
        
        # Parse the XML
        new_project = xml_parser.parse_string(xml_content)
        
        logger.info(f"Parsed project: {new_project.project_name} ({new_project.project_code})")
        logger.info(f"Milestones found: {len(new_project.milestones)}")
        logger.info(f"Risks found: {len(new_project.risks)}")
        logger.info(f"Changes found: {len(new_project.changes)}")
        
        # Debug: Log first few milestones if any exist
        if new_project.milestones:
            for i, milestone in enumerate(new_project.milestones[:3]):
                logger.info(f"Milestone {i+1}: {milestone.name} - {milestone.target_date} - Parent: {milestone.parent_project}")
        else:
            logger.warning("No milestones extracted from XML!")
        
        # Check if project already exists
        existing_project = project_repo.get_project_by_code(
            new_project.project_code
        )
        
        detected_changes = []
        
        # Only detect changes if NOT a baseline upload
        if existing_project and not is_baseline_upload:
            logger.info("🔍 Checking for milestone date changes...")
            logger.info(
                f"   Old: {len(existing_project.milestones)} milestones"
            )
            logger.info(
                f"   New: {len(new_project.milestones)} milestones"
            )
            
            # Detect changes
            changes = change_detector.detect_milestone_changes(
                existing_project,
                new_project
            )
            
            logger.info(f"📊 DETECTED {len(changes)} DATE CHANGES")
            
            if changes:
                for c in changes:
                    logger.info(
                        f"   • {c['milestone_name']}: "
                        f"{c['old_date']} → {c['new_date']} "
                        f"({c['days_diff']:+d} days)"
                    )
            
            detected_changes = [
                {
                    'milestone_name': c['milestone_name'],
                    'old_date': c['old_date'],
                    'new_date': c['new_date'],
                    'days_diff': c['days_diff'],
                    'type': c['type'],
                    'suggested_impact': change_detector.calculate_impact(
                        c['days_diff'],
                        c['milestone_name']
                    ),
                    'has_reason': c['existing_change'] is not None,
                    'existing_reason': (
                        c['existing_change'].reason 
                        if c['existing_change'] else ''
                    ),
                    'existing_impact': (
                        c['existing_change'].impact 
                        if c['existing_change'] else ''
                    )
                }
                for c in changes
            ]
        
        # Save uploaded file for reference
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{new_project.project_code}_{timestamp}.xml"
        upload_path = UPLOAD_DIR / filename
        upload_path.write_text(xml_content)
        
        logger.info(f"Saved upload to {upload_path}")
        
        # **FIX: Always save the project data, not just when confirming changes**
        # This ensures milestones, gantt data, and risks are updated
        project_dir = DATA_DIR / f"PROJECT-{new_project.project_code.replace('-', '_')}"
        project_dir.mkdir(exist_ok=True)
        
        yaml_path = project_dir / "project_status.yaml"
        
        # Merge with existing changes if this is an update (not baseline)
        clear_old_changes = clear_previous_changes.lower() == "true"
        
        if existing_project and not is_baseline_upload:
            # AUTO-SAVE detected changes with auto-generated reason
            auto_changes = []
            for c in detected_changes:
                change = change_detector.create_change_record(
                    change_info={
                        'milestone_name': c['milestone_name'],
                        'old_date': c['old_date'],
                        'new_date': c['new_date'],
                        'days_diff': c['days_diff']
                    },
                    reason=f"Schedule change detected on {datetime.now().strftime('%Y-%m-%d')}",
                    impact=c['suggested_impact'],
                    project_code=new_project.project_code
                )
                auto_changes.append(change)
            
            # Conditionally merge or replace changes
            if clear_old_changes:
                # Replace all old changes with only new detected changes
                logger.info(f"🗑️ Clearing {len(existing_project.changes)} previous changes")
                new_project.changes = auto_changes
                logger.info(f"✅ Keeping only {len(auto_changes)} new changes")
            else:
                # Merge auto-detected changes with existing changes
                new_project.changes = change_detector.merge_changes(
                    existing_project.changes,
                    auto_changes
                )
                logger.info(f"📊 Merged: {len(existing_project.changes)} old + {len(auto_changes)} new = {len(new_project.changes)} total")
        
        # Save risks from XML - merge with existing if updating
        yaml_risks = new_project.risks  # Start with risks from new XML
        
        # If updating (not baseline), preserve existing risks that aren't in the new XML
        if existing_project and hasattr(existing_project, 'risks') and not is_baseline_upload:
            # Keep existing risks, update with new ones from XML by risk_id
            existing_risk_ids = {r.risk_id for r in yaml_risks}
            for existing_risk in existing_project.risks:
                if existing_risk.risk_id not in existing_risk_ids:
                    yaml_risks.append(existing_risk)
        
        # SMART MERGE: Preserve manually edited milestone data
        milestones_to_save = []
        seen_milestone_names = set()  # Track to prevent duplicates
        
        if existing_project and not is_baseline_upload and hasattr(
            existing_project, 'milestones'
        ):
            # Create lookup by milestone ID and name
            existing_by_id = {
                getattr(m, 'id', None): m
                for m in existing_project.milestones
                if getattr(m, 'id', None)
            }
            existing_by_name = {
                m.name: m for m in existing_project.milestones
            }
            
            for new_milestone in new_project.milestones:
                new_id = getattr(new_milestone, 'id', None)
                milestone_name = new_milestone.name.strip()
                
                # Skip if already processed (duplicate in XML)
                if milestone_name in seen_milestone_names:
                    logger.warning(
                        f"⚠️ DUPLICATE in XML: '{milestone_name}' "
                        f"at line {new_project.milestones.index(new_milestone)} "
                        f"- skipping duplicate"
                    )
                    continue
                seen_milestone_names.add(milestone_name)
                
                # Try to find existing milestone with multiple strategies
                existing = None
                match_method = None
                
                # Strategy 1: Match by ID (most reliable)
                if new_id and new_id in existing_by_id:
                    existing = existing_by_id[new_id]
                    match_method = 'id'
                
                # Strategy 2: Match by name
                elif milestone_name in existing_by_name:
                    existing = existing_by_name[milestone_name]
                    match_method = 'name'
                
                # Strategy 3: Match by date + parent (handles renamed milestones)
                else:
                    for existing_m in existing_project.milestones:
                        if (existing_m.target_date == new_milestone.target_date and
                            existing_m.parent_project == new_milestone.parent_project and
                            new_milestone.target_date and new_milestone.parent_project):
                            existing = existing_m
                            match_method = 'date+parent'
                            logger.warning(
                                f"📝 Matched by date+parent: "
                                f"'{existing_m.name}' ← '{milestone_name}'"
                            )
                            break
                
                if existing:
                    # Log the match for debugging
                    if match_method == 'date+parent':
                        logger.info(
                            f"Merging renamed milestone: "
                            f"XML '{milestone_name}' → "
                            f"Existing '{existing.name}'"
                        )
                    
                    # SMART NAME HANDLING:
                    # - If matched by ID or date+parent, use XML name (handles MS Project renames)
                    # - If matched by name, preserve existing name (protects manual edits)
                    use_xml_name = match_method in ['id', 'date+parent']
                    final_name = milestone_name if use_xml_name else existing.name
                    
                    if use_xml_name and milestone_name != existing.name:
                        logger.info(
                            f"📝 Updating milestone name: "
                            f"'{existing.name}' → '{milestone_name}'"
                        )
                    
                    # PRESERVE user edits for certain fields
                    # REFRESH from XML for MS Project data
                    merged_milestone = {
                        'id': new_id,
                        'name': final_name,  # Use XML name for renamed milestones
                        'target_date': new_milestone.target_date,  # XML
                        'status': new_milestone.status,  # XML
                        'completion_date': new_milestone.completion_date,  # XML
                        'completion_percentage': new_milestone.completion_percentage,  # XML
                        'notes': existing.notes,  # PRESERVE user edits
                        'parent_project': new_milestone.parent_project,  # XML
                        'resources': existing.resources,  # PRESERVE user edits
                        'project': new_project.project_code
                    }
                    milestones_to_save.append(merged_milestone)
                else:
                    # New milestone not in existing data - add it
                    milestones_to_save.append({
                        'id': new_id,
                        'name': new_milestone.name,
                        'target_date': new_milestone.target_date,
                        'status': new_milestone.status,
                        'completion_date': new_milestone.completion_date,
                        'completion_percentage': new_milestone.completion_percentage,
                        'notes': new_milestone.notes,
                        'parent_project': new_milestone.parent_project,
                        'resources': new_milestone.resources,
                        'project': new_project.project_code
                    })
        else:
            # First upload or baseline - deduplicate and use milestones from XML
            seen_names = set()
            for m in new_project.milestones:
                milestone_name = m.name.strip()
                if milestone_name in seen_names:
                    logger.warning(
                        f"⚠️ DUPLICATE MILESTONE IN XML: '{milestone_name}' "
                        f"- skipping duplicate"
                    )
                    continue
                seen_names.add(milestone_name)
                
                milestones_to_save.append({
                    'id': getattr(m, 'id', None),
                    'name': m.name,
                    'target_date': m.target_date,
                    'status': m.status,
                    'completion_date': m.completion_date,
                    'completion_percentage': m.completion_percentage,
                    'notes': m.notes,
                    'parent_project': m.parent_project,
                    'resources': m.resources,
                    'project': new_project.project_code
                })
        
        # CLEANUP: Remove any existing duplicates from final merged data
        final_milestones = []
        final_seen_names = set()
        duplicates_removed = 0
        
        for m in milestones_to_save:
            m_name = m.get('name', '').strip()
            if m_name in final_seen_names:
                duplicates_removed += 1
                logger.warning(
                    f"⚠️ REMOVING DUPLICATE from saved data: '{m_name}'"
                )
                continue
            final_seen_names.add(m_name)
            final_milestones.append(m)
        
        if duplicates_removed > 0:
            logger.warning(
                f"🧹 Cleaned up {duplicates_removed} duplicate milestone(s)"
            )
        
        milestones_to_save = final_milestones
        
        # Convert to dict for YAML serialization
        project_dict = {
            'project_name': new_project.project_name,
            'project_code': new_project.project_code,
            'status': new_project.status,
            'start_date': new_project.start_date,
            'target_completion': new_project.target_completion,
            'completion_percentage': new_project.completion_percentage,
            'milestones': milestones_to_save,
            # Save risks from XML to project YAML
            'risks': [
                {
                    'risk_id': r.risk_id,
                    'description': r.description,
                    'severity': r.severity,
                    'probability': r.probability,
                    'impact': r.impact,
                    'mitigation': r.mitigation,
                    'status': r.status
                }
                for r in yaml_risks
            ],
            'changes': [
                {
                    'change_id': c.change_id,
                    'date': c.date,
                    'old_date': c.old_date,
                    'new_date': c.new_date,
                    'reason': c.reason,
                    'impact': c.impact
                }
                for c in new_project.changes
            ]
        }
        
        with open(yaml_path, 'w') as f:
            yaml.dump(project_dict, f, default_flow_style=False, sort_keys=False)
        
        logger.info(f"Project data saved to {yaml_path}")
        logger.info(f"Saved {len(new_project.milestones)} milestones, {len(new_project.risks)} risks")
        
        # Record the upload for subscription tracking
        try:
            sub_service.record_project_upload(
                user_id=user.user_id,
                project_name=new_project.project_name,
                file_size_mb=file_size_mb,
                xml_file_path=str(upload_path)
            )
            logger.info("Recorded upload in subscription tracking")
        except Exception as track_error:
            # Don't fail the upload if tracking fails
            logger.error(f"Failed to record upload in subscription tracking: {track_error}")
        
        logger.info("Upload completed successfully")
        
        return JSONResponse({
            'success': True,
            'project_code': new_project.project_code,
            'project_name': new_project.project_name,
            'is_new': existing_project is None,
            'detected_changes': detected_changes,
            'milestone_count': len(new_project.milestones),
            'upload_path': str(upload_path)
        })
        
    except SubscriptionError as se:
        logger.error(f"Subscription error: {se.detail}")
        return JSONResponse({
            'success': False,
            'error': se.detail,
            'error_type': 'subscription_limit'
        }, status_code=402)
    except Exception as e:
        logger.error(f"Upload failed with error: {str(e)}", exc_info=True)
        return JSONResponse({
            'success': False,
            'error': str(e),
            'error_type': 'general_error'
        }, status_code=400)


@router.post("/upload/confirm")
async def confirm_upload(
    project_code: str = Form(...),
    upload_path: str = Form(...),
    changes_json: str = Form(...)
):
    """
    Confirm upload and save project with change reasons
    """
    import json
    import logging
    logger = logging.getLogger(__name__)
    
    try:
        logger.info(f"Confirm upload called for project: {project_code}")
        logger.info(f"Upload path: {upload_path}")
        logger.info(f"Changes JSON: {changes_json[:200]}...")
        
        # Verify upload path exists
        upload_file = Path(upload_path)
        if not upload_file.exists():
            logger.error(f"Upload file not found: {upload_path}")
            raise FileNotFoundError(f"Upload file not found: {upload_path}")
        
        # Read the uploaded XML again
        xml_content = upload_file.read_text()
        logger.info(
            f"Read XML file successfully, size: {len(xml_content)} bytes"
        )
        
        new_project = xml_parser.parse_string(xml_content)
        logger.info(f"Parsed project: {new_project.project_name}")
        
        # Parse changes from form
        changes_data = json.loads(changes_json)
        
        # Get existing project to merge changes
        existing_project = project_repo.get_project_by_code(project_code)
        
        # Create Change objects from submitted data
        new_changes = []
        for change_data in changes_data:
            if change_data.get('reason'):  # Only if reason provided
                change = change_detector.create_change_record(
                    change_info={
                        'milestone_name': change_data['milestone_name'],
                        'old_date': change_data['old_date'],
                        'new_date': change_data['new_date'],
                        'days_diff': change_data['days_diff']
                    },
                    reason=change_data['reason'],
                    impact=change_data.get('impact', 'Not specified'),
                    project_code=project_code
                )
                new_changes.append(change)
        
        # Merge with existing changes
        if existing_project:
            all_changes = change_detector.merge_changes(
                existing_project.changes,
                new_changes
            )
            new_project.changes = all_changes
        else:
            new_project.changes = new_changes
        
        # Save to YAML file
        project_dir = DATA_DIR / f"PROJECT-{project_code.replace('-', '_')}"
        project_dir.mkdir(exist_ok=True)
        
        yaml_path = project_dir / "project_status.yaml"
        
        # Convert to dict for YAML serialization
        project_dict = {
            'project_name': new_project.project_name,
            'project_code': new_project.project_code,
            'status': new_project.status,
            'start_date': new_project.start_date,
            'target_completion': new_project.target_completion,
            'completion_percentage': new_project.completion_percentage,
            'milestones': [
                {
                    'name': m.name,
                    'target_date': m.target_date,
                    'status': m.status,
                    'completion_date': m.completion_date,
                    'completion_percentage': m.completion_percentage,
                    'notes': m.notes,
                    'parent_project': m.parent_project,
                    'resources': m.resources,
                    'project': new_project.project_code  # Add project code for frontend
                }
                for m in new_project.milestones
            ],
            'risks': [
                {
                    'risk_id': r.risk_id,
                    'description': r.description,
                    'severity': r.severity,
                    'probability': r.probability,
                    'impact': r.impact,
                    'mitigation': r.mitigation,
                    'status': r.status
                }
                for r in new_project.risks
            ],
            'changes': [
                {
                    'change_id': c.change_id,
                    'date': c.date,
                    'old_date': c.old_date,
                    'new_date': c.new_date,
                    'reason': c.reason,
                    'impact': c.impact
                }
                for c in new_project.changes
            ]
        }
        
        with open(yaml_path, 'w') as f:
            yaml.dump(project_dict, f, default_flow_style=False, sort_keys=False)
        
        # Log first milestone to verify data
        if new_project.milestones:
            first_m = new_project.milestones[0]
            logger.info(f"First milestone data - Name: {first_m.name}, "
                       f"Parent: {first_m.parent_project}, "
                       f"Resources: {first_m.resources}")
        
        logger.info(f"Project {project_code} saved successfully")
        
        return JSONResponse({
            'success': True,
            'message': f'Project {project_code} updated successfully',
            'changes_saved': len(new_changes)
        })
        
    except Exception as e:
        logger.error(f"Confirm upload failed: {str(e)}", exc_info=True)
        return JSONResponse({
            'success': False,
            'error': str(e)
        }, status_code=400)


@router.post("/changes/{project_code}/update")
async def update_change_reason(
    project_code: str,
    change_id: str = Form(...),
    reason: str = Form(...),
    impact: str = Form(...)
):
    """Update change reason inline"""
    try:
        # Load project
        project = project_repo.get_project_by_code(project_code)
        
        if not project:
            return JSONResponse({
                'success': False,
                'error': 'Project not found'
            }, status_code=404)
        
        # Find and update change
        updated = False
        for change in project.changes:
            if change.change_id == change_id:
                change.reason = reason
                change.impact = impact
                updated = True
                break
        
        if not updated:
            return JSONResponse({
                'success': False,
                'error': 'Change not found'
            }, status_code=404)
        
        # Save back to YAML
        project_dir = DATA_DIR / f"PROJECT-{project_code.replace('-', '_')}"
        yaml_path = project_dir / "project_status.yaml"
        
        project_dict = {
            'project_name': project.project_name,
            'project_code': project.project_code,
            'status': project.status,
            'start_date': project.start_date,
            'target_completion': project.target_completion,
            'completion_percentage': project.completion_percentage,
            'milestones': [
                {
                    'name': m.name,
                    'target_date': m.target_date,
                    'status': m.status,
                    'completion_date': m.completion_date,
                    'completion_percentage': m.completion_percentage,
                    'notes': m.notes,
                    'parent_project': getattr(m, 'parent_project', None),
                    'resources': getattr(m, 'resources', None),
                    'project': project.project_code  # Add project code for frontend
                }
                for m in project.milestones
            ],
            'risks': [
                {
                    'risk_id': r.risk_id,
                    'description': r.description,
                    'severity': r.severity,
                    'probability': r.probability,
                    'impact': r.impact,
                    'mitigation': r.mitigation,
                    'status': r.status
                }
                for r in project.risks
            ],
            'changes': [
                {
                    'change_id': c.change_id,
                    'date': c.date,
                    'old_date': c.old_date,
                    'new_date': c.new_date,
                    'reason': c.reason,
                    'impact': c.impact
                }
                for c in project.changes
            ]
        }
        
        with open(yaml_path, 'w') as f:
            yaml.dump(project_dict, f, default_flow_style=False, sort_keys=False)
        
        return JSONResponse({
            'success': True,
            'message': 'Change updated successfully'
        })
        
    except Exception as e:
        return JSONResponse({
            'success': False,
            'error': str(e)
        }, status_code=400)


@router.post("/clear-cache/{project_code}")
async def clear_project_cache(project_code: str):
    """Clear cached project data - force fresh baseline on next upload"""
    try:
        import logging
        logger = logging.getLogger(__name__)
        
        logger.info(f"🗑️ Clearing cache for project: {project_code}")
        
        # Delete the project YAML file
        project_dir = DATA_DIR / f"PROJECT-{project_code.replace('-', '_')}"
        yaml_path = project_dir / "project_status.yaml"
        
        if yaml_path.exists():
            yaml_path.unlink()
            logger.info(f"✅ Deleted: {yaml_path}")
            
            # Also clear from in-memory cache
            project_repo._projects = [
                p for p in project_repo._projects
                if p.project_code != project_code
            ]
            
            return JSONResponse({
                'success': True,
                'message': (
                    f'Cache cleared for {project_code}. '
                    'Next upload will be treated as baseline.'
                )
            })
        else:
            return JSONResponse({
                'success': False,
                'message': f'No cached data found for {project_code}'
            }, status_code=404)
            
    except Exception as e:
        logger.error(f"❌ Failed to clear cache: {e}")
        return JSONResponse({
            'success': False,
            'error': str(e)
        }, status_code=500)


# ===== PowerPoint Template Management =====

@router.post("/upload/powerpoint-template")
async def upload_powerpoint_template(
    file: UploadFile = File(...),
    name: str = Form(...)
):
    """Upload a PowerPoint template with organization branding"""
    try:
        # Validate file type
        if not file.filename.endswith('.pptx'):
            return JSONResponse({
                'success': False,
                'detail': 'Only .pptx PowerPoint files are accepted'
            }, status_code=400)
        
        # Generate unique template ID
        template_id = f"template_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        template_path = POWERPOINT_TEMPLATES_DIR / f"{template_id}.pptx"
        
        # Save file
        content = await file.read()
        with open(template_path, 'wb') as f:
            f.write(content)
        
        # Save metadata
        metadata = {
            'id': template_id,
            'name': name,
            'filename': file.filename,
            'uploaded_at': datetime.now().isoformat(),
            'is_default': False
        }
        
        metadata_path = POWERPOINT_TEMPLATES_DIR / f"{template_id}.yaml"
        with open(metadata_path, 'w') as f:
            yaml.dump(metadata, f)
        
        return JSONResponse({
            'success': True,
            'template_id': template_id,
            'name': name
        })
        
    except Exception as e:
        logger.error(f"❌ Failed to upload template: {e}")
        return JSONResponse({
            'success': False,
            'detail': str(e)
        }, status_code=500)


@router.get("/upload/powerpoint-templates")
async def list_powerpoint_templates():
    """List all available PowerPoint templates"""
    try:
        templates = []
        
        # Find all template metadata files
        for metadata_file in POWERPOINT_TEMPLATES_DIR.glob("template_*.yaml"):
            try:
                with open(metadata_file, 'r') as f:
                    metadata = yaml.safe_load(f)
                    
                # Check if template file exists
                template_path = POWERPOINT_TEMPLATES_DIR / f"{metadata['id']}.pptx"
                if template_path.exists():
                    templates.append(metadata)
            except Exception as e:
                logger.warning(f"Failed to load template metadata {metadata_file}: {e}")
        
        # Sort by upload date (newest first)
        templates.sort(key=lambda x: x.get('uploaded_at', ''), reverse=True)
        
        return templates
        
    except Exception as e:
        logger.error(f"❌ Failed to list templates: {e}")
        return JSONResponse({
            'success': False,
            'detail': str(e)
        }, status_code=500)


@router.delete("/upload/powerpoint-template/{template_id}")
async def delete_powerpoint_template(template_id: str):
    """Delete a PowerPoint template"""
    try:
        template_path = POWERPOINT_TEMPLATES_DIR / f"{template_id}.pptx"
        metadata_path = POWERPOINT_TEMPLATES_DIR / f"{template_id}.yaml"
        
        if not template_path.exists():
            return JSONResponse({
                'success': False,
                'detail': 'Template not found'
            }, status_code=404)
        
        # Delete files
        template_path.unlink()
        if metadata_path.exists():
            metadata_path.unlink()
        
        # Also delete cached preview
        cache_dir = POWERPOINT_TEMPLATES_DIR / "previews"
        for cached_file in cache_dir.glob(f"{template_id}_*.png"):
            cached_file.unlink()
        
        return JSONResponse({
            'success': True,
            'message': f'Template {template_id} deleted'
        })
        
    except Exception as e:
        logger.error(f"❌ Failed to delete template: {e}")
        return JSONResponse({
            'success': False,
            'detail': str(e)
        }, status_code=500)


@router.delete("/upload/powerpoint-template/{template_id}/preview-cache")
async def clear_template_preview_cache(template_id: str):
    """Clear the cached preview for a template to force regeneration"""
    try:
        cache_dir = POWERPOINT_TEMPLATES_DIR / "previews"
        deleted_count = 0
        for cached_file in cache_dir.glob(f"{template_id}_*.png"):
            cached_file.unlink()
            deleted_count += 1
        
        return JSONResponse({
            'success': True,
            'message': f'Cleared {deleted_count} cached preview(s) for {template_id}'
        })
        
    except Exception as e:
        logger.error(f"❌ Failed to clear cache: {e}")
        return JSONResponse({
            'success': False,
            'detail': str(e)
        }, status_code=500)


@router.post("/upload/powerpoint-template/{template_id}/set-default")
async def set_default_template(template_id: str):
    """Set a template as the default for PowerPoint exports"""
    try:
        # Clear all existing defaults
        for metadata_file in POWERPOINT_TEMPLATES_DIR.glob("template_*.yaml"):
            try:
                with open(metadata_file, 'r') as f:
                    metadata = yaml.safe_load(f)
                metadata['is_default'] = False
                with open(metadata_file, 'w') as f:
                    yaml.dump(metadata, f)
            except Exception as e:
                logger.warning(f"Failed to update metadata {metadata_file}: {e}")
        
        # Set new default
        metadata_path = POWERPOINT_TEMPLATES_DIR / f"{template_id}.yaml"
        if not metadata_path.exists():
            return JSONResponse({
                'success': False,
                'detail': 'Template not found'
            }, status_code=404)
        
        with open(metadata_path, 'r') as f:
            metadata = yaml.safe_load(f)
        metadata['is_default'] = True
        with open(metadata_path, 'w') as f:
            yaml.dump(metadata, f)
        
        return JSONResponse({
            'success': True,
            'message': f'Template {template_id} set as default'
        })
        
    except Exception as e:
        logger.error(f"❌ Failed to set default template: {e}")
        return JSONResponse({
            'success': False,
            'detail': str(e)
        }, status_code=500)


@router.get("/upload/powerpoint-template/{template_id}/preview")
async def get_template_preview(template_id: str, slide_index: int = 0):
    """Get a preview image of a template slide.
    
    Dynamically renders all visual elements from the template including:
    - All images (logos, graphics) at their actual positions
    - Shape outlines and fills
    - Placeholder areas (title, content, footer)
    - Lines and separators
    
    This works with any template layout.
    """
    from fastapi.responses import FileResponse
    from pptx import Presentation
    from pptx.util import Emu, Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from PIL import Image, ImageDraw, ImageFont
    from io import BytesIO
    import zipfile
    import re
    
    try:
        template_path = POWERPOINT_TEMPLATES_DIR / f"{template_id}.pptx"
        metadata_path = POWERPOINT_TEMPLATES_DIR / f"{template_id}.yaml"
        
        if not template_path.exists():
            return JSONResponse({
                'success': False,
                'detail': 'Template not found'
            }, status_code=404)
        
        # Load metadata for template name
        template_name = template_id
        if metadata_path.exists():
            with open(metadata_path, 'r') as f:
                metadata = yaml.safe_load(f)
                template_name = metadata.get('name', template_id)
        
        # Check if cached preview exists (v4 = full dynamic rendering)
        cache_dir = POWERPOINT_TEMPLATES_DIR / "previews"
        cache_dir.mkdir(exist_ok=True)
        cached_preview = cache_dir / f"{template_id}_slide{slide_index}_v4.png"
        
        if cached_preview.exists():
            return FileResponse(cached_preview, media_type="image/png")
        
        # Create template preview by rendering all elements
        try:
            prs = Presentation(template_path)
            
            # Scale factors: slide EMU to 1920x1080 pixels
            px_per_emu_x = 1920 / prs.slide_width
            px_per_emu_y = 1080 / prs.slide_height
            
            def to_px(left, top, width, height):
                return (
                    int(left * px_per_emu_x),
                    int(top * px_per_emu_y),
                    int(width * px_per_emu_x),
                    int(height * px_per_emu_y)
                )
            
            # Create white base image
            img = Image.new('RGBA', (1920, 1080), color=(255, 255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # Load all images from the PPTX media folder
            media_images = {}
            with zipfile.ZipFile(template_path, 'r') as zf:
                for name in zf.namelist():
                    if name.startswith('ppt/media/') and any(name.endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif']):
                        try:
                            with zf.open(name) as f:
                                media_images[name.split('/')[-1]] = Image.open(BytesIO(f.read())).copy()
                        except:
                            pass
            
            logger.info(f"Loaded {len(media_images)} media images from template")
            
            # Get fonts
            try:
                font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
                font_normal = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
                font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
            except:
                font_title = font_normal = font_small = ImageFont.load_default()
            
            # Function to render shapes from a shape collection
            def render_shapes(shapes, is_master=False):
                for shape in shapes:
                    left_px, top_px, width_px, height_px = to_px(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    
                    # Skip off-screen elements
                    if left_px > 1920 or top_px > 1080:
                        continue
                    if width_px < 2 or height_px < 2:
                        continue
                    
                    shape_type = shape.shape_type
                    name = shape.name.lower()
                    
                    # Render PICTURES (logos, graphics)
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Try to get the image blob
                            if hasattr(shape, 'image') and shape.image:
                                pic_data = shape.image.blob
                                pic_img = Image.open(BytesIO(pic_data))
                                
                                # Resize to match shape dimensions
                                pic_img = pic_img.resize((max(1, width_px), max(1, height_px)), Image.Resampling.LANCZOS)
                                
                                if pic_img.mode != 'RGBA':
                                    pic_img = pic_img.convert('RGBA')
                                
                                # Paste with alpha if available
                                img.paste(pic_img, (left_px, top_px), pic_img)
                                logger.info(f"Rendered picture: {shape.name} at ({left_px}, {top_px})")
                        except Exception as e:
                            # Draw placeholder for failed images
                            draw.rectangle([(left_px, top_px), (left_px + width_px, top_px + height_px)],
                                          outline=(200, 200, 210), width=1)
                    
                    # Render LINES
                    elif shape_type == MSO_SHAPE_TYPE.LINE:
                        draw.line([(left_px, top_px), (left_px + width_px, top_px + height_px)],
                                 fill=(180, 185, 195), width=2)
                    
                    # Render AUTO_SHAPES (rectangles, etc.)
                    elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        # Check for fill
                        fill_color = None
                        try:
                            if hasattr(shape, 'fill') and shape.fill.type is not None:
                                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                                    rgb = shape.fill.fore_color.rgb
                                    if rgb:
                                        fill_color = (rgb[0], rgb[1], rgb[2], 50)  # Semi-transparent
                        except:
                            pass
                        
                        if fill_color:
                            overlay = Image.new('RGBA', (width_px, height_px), fill_color)
                            img.paste(overlay, (left_px, top_px), overlay)
                        
                        # Draw outline
                        draw.rectangle([(left_px, top_px), (left_px + width_px, top_px + height_px)],
                                      outline=(210, 215, 225), width=1)
                    
                    # Render PLACEHOLDERS (title, content, footer areas)
                    elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Determine placeholder type from name
                        is_title = 'titre' in name or 'title' in name
                        is_content = 'texte' in name or 'content' in name or 'text' in name
                        is_footer = 'pied' in name or 'footer' in name
                        
                        if is_title:
                            # Title area - light blue tint
                            overlay = Image.new('RGBA', (width_px, height_px), (230, 240, 255, 80))
                            img.paste(overlay, (left_px, top_px), overlay)
                            draw.rectangle([(left_px, top_px), (left_px + width_px, top_px + height_px)],
                                          outline=(150, 180, 220), width=2)
                            draw.text((left_px + 10, top_px + 5), "TITLE AREA", fill=(120, 150, 190), font=font_small)
                        
                        elif is_content and height_px > 200:
                            # Content area - dashed border
                            for i in range(left_px, left_px + width_px, 15):
                                draw.line([(i, top_px), (min(i+8, left_px + width_px), top_px)], fill=(200, 210, 225), width=1)
                                draw.line([(i, top_px + height_px), (min(i+8, left_px + width_px), top_px + height_px)], fill=(200, 210, 225), width=1)
                            for i in range(top_px, top_px + height_px, 15):
                                draw.line([(left_px, i), (left_px, min(i+8, top_px + height_px))], fill=(200, 210, 225), width=1)
                                draw.line([(left_px + width_px, i), (left_px + width_px, min(i+8, top_px + height_px))], fill=(200, 210, 225), width=1)
                            # Label
                            label = "CONTENT AREA - Place your screenshot here"
                            bbox = draw.textbbox((0, 0), label, font=font_normal)
                            lx = left_px + (width_px - (bbox[2] - bbox[0])) // 2
                            ly = top_px + (height_px - (bbox[3] - bbox[1])) // 2
                            draw.text((lx, ly), label, fill=(180, 190, 210), font=font_normal)
                        
                        elif is_footer:
                            # Footer area - subtle gray
                            overlay = Image.new('RGBA', (width_px, height_px), (245, 245, 250, 100))
                            img.paste(overlay, (left_px, top_px), overlay)
                            draw.line([(left_px, top_px), (left_px + width_px, top_px)], fill=(220, 225, 235), width=1)
                    
                    # Render GROUPS (may contain multiple elements)
                    elif shape_type == MSO_SHAPE_TYPE.GROUP:
                        # Draw group boundary
                        draw.rectangle([(left_px, top_px), (left_px + width_px, top_px + height_px)],
                                      outline=(190, 200, 215), width=1)
                    
                    # Render TEXT_BOX
                    elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        if hasattr(shape, 'text') and shape.text.strip():
                            # Small text boxes with content (like confidentiality marks)
                            draw.rectangle([(left_px, top_px), (left_px + width_px, top_px + height_px)],
                                          fill=(250, 250, 252), outline=(230, 235, 245), width=1)
            
            # Render slide master elements (persistent across all slides)
            logger.info("Rendering slide master elements...")
            render_shapes(prs.slide_master.shapes, is_master=True)
            
            # Render the specified slide layout (or first content layout)
            layout_idx = min(slide_index, len(prs.slide_layouts) - 1)
            # Try to find a good content layout (skip cover pages)
            for i, layout in enumerate(prs.slide_layouts):
                if any(x in layout.name.lower() for x in ['content', 'text', 'title_text', 'white']):
                    layout_idx = i
                    break
            
            if layout_idx < len(prs.slide_layouts):
                logger.info(f"Rendering slide layout: {prs.slide_layouts[layout_idx].name}")
                render_shapes(prs.slide_layouts[layout_idx].shapes)
            
            # Add template name at bottom
            draw.text((20, 1055), f"Template: {template_name}", fill=(150, 155, 170), font=font_small)
            
            # Convert to RGB for PNG saving
            img = img.convert('RGB')
            
            # Save to cache
            img.save(cached_preview, 'PNG')
            logger.info(f"Generated template preview v4: {cached_preview}")
            
            return FileResponse(cached_preview, media_type="image/png")
            
        except Exception as e:
            logger.error(f"Failed to generate preview from PPTX: {e}")
            # Fall through to placeholder
        
        # Fallback: Generate a styled placeholder
        img = Image.new('RGB', (1920, 1080), color=(250, 250, 252))
        draw = ImageDraw.Draw(img)
        
        try:
            font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
            font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except:
            font_large = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        # Draw a subtle header area
        draw.rectangle([(0, 0), (1920, 80)], fill=(245, 245, 250))
        draw.line([(0, 80), (1920, 80)], fill=(220, 220, 225), width=1)
        
        # Draw template name in header
        draw.text((30, 25), template_name, fill=(60, 60, 70), font=font_large)
        
        # Draw content area placeholder with dashed border
        for i in range(0, 1840, 20):
            draw.line([(40 + i, 100), (min(50 + i, 1880), 100)], fill=(200, 200, 210), width=1)
            draw.line([(40 + i, 1040), (min(50 + i, 1880), 1040)], fill=(200, 200, 210), width=1)
        
        # Draw "Content Area" text
        text = "Drag your content here"
        bbox = draw.textbbox((0, 0), text, font=font_small)
        x = (1920 - (bbox[2] - bbox[0])) // 2
        draw.text((x, 560), text, fill=(180, 180, 180), font=font_small)
        
        # Save to cache
        img.save(cached_preview, 'PNG')
        
        return FileResponse(cached_preview, media_type="image/png")
        
    except Exception as e:
        logger.error(f"❌ Failed to get template preview: {e}")
        return JSONResponse({
            'success': False,
            'detail': str(e)
        }, status_code=500)




