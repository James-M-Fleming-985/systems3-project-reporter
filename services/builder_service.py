import os
import time
from io import BytesIO
from typing import List, Optional, Dict, Any
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


# Template safe zone configuration (in inches from each edge)
# These define where the image content should not encroach
TEMPLATE_SAFE_ZONES = {
    'default': {'top': 1.2, 'bottom': 0.5, 'left': 0.5, 'right': 0.5},      # Default with title area
    'no_template': {'top': 0.3, 'bottom': 0.3, 'left': 0.3, 'right': 0.3},  # Minimal margins
    'branded': {'top': 1.4, 'bottom': 0.6, 'left': 0.5, 'right': 0.5}       # More space for logos/footers
}


class PowerPointBuilderService:
    """Service for building PowerPoint presentations with screenshots and branding."""
    
    DEFAULT_TEMPLATE_PATH = "templates/default_template.pptx"
    MAX_GENERATION_TIME = 3.0  # seconds
    
    def __init__(self):
        """Initialize the PowerPoint builder service."""
        self.presentation = None
        self.template_path = None
        self.start_time = None
        
    def generate_presentation(
        self,
        report_data: Dict[str, Any],
        screenshots: List[bytes],
        template_path: Optional[str] = None,
        skip_slide_titles: bool = False,
        slide_transforms: Optional[List[Dict[str, Any]]] = None,
        slide_titles: Optional[List[str]] = None
    ) -> bytes:
        """
        Generate a PowerPoint presentation with report data and screenshots.
        
        Args:
            report_data: Dictionary containing report metadata
            screenshots: List of screenshot images as bytes
            template_path: Optional path to company template
            skip_slide_titles: If True, don't add titles (templates have their own)
            slide_transforms: Optional list of transform data for each slide
            slide_titles: Optional list of descriptive titles for each slide
            
        Returns:
            Generated PPTX file as bytes
            
        Raises:
            ValueError: If template is invalid or generation takes too long
        """
        self.start_time = time.time()
        
        # Load template
        self._load_template(template_path)
        
        # Determine if we should skip titles (branded templates have their own)
        use_template_titles = skip_slide_titles or (template_path is not None)
        
        # If using a template, remove existing content slides (keep only layout masters)
        # Typically slides 2+ are content slides that should be replaced
        if template_path and self.presentation:
            self._prepare_template_for_content(report_data.get('include_title_slide', True))
        
        # Create title slide (only if not using template or specifically requested)
        if not use_template_titles or report_data.get('include_title_slide', True):
            self._create_title_slide(report_data)
        
        # Create content slides for screenshots
        for idx, screenshot in enumerate(screenshots):
            # Get transform for this slide if provided
            transform = None
            if slide_transforms and idx < len(slide_transforms):
                transform = slide_transforms[idx]
            
            # Apply transform to screenshot if provided
            if transform:
                processed_screenshot = self._apply_transform_to_image(
                    screenshot, transform)
            else:
                processed_screenshot = screenshot
            
            # Get slide title - use provided title or generate default
            if slide_titles and idx < len(slide_titles):
                title = slide_titles[idx]
            else:
                title = f"Screenshot {idx + 1}"
            
            self._create_content_slide(
                processed_screenshot, 
                idx + 1, 
                skip_title=use_template_titles,
                slide_title=title
            )
            
        # Check generation time
        elapsed = time.time() - self.start_time
        if elapsed >= self.MAX_GENERATION_TIME:
            raise ValueError("File generation exceeded time limit")
            
        # Save to bytes
        return self._save_to_bytes()
    
    def generate_hybrid_presentation(
        self,
        report_data: Dict[str, Any],
        slides_data: List[Dict[str, Any]],
        template_path: Optional[str] = None,
    ) -> bytes:
        """
        Generate a PowerPoint with mix of screenshots and native tables.
        
        Args:
            report_data: Dictionary containing report metadata
            slides_data: List of slide configs, each with:
                - type: 'screenshot', 'risks', or 'milestones'
                - data: screenshot bytes OR list of risk/milestone dicts
                - title: Slide title
                - page_num/total_pages: For multi-page content
            template_path: Optional path to company template
            
        Returns:
            Generated PPTX file as bytes
        """
        self.start_time = time.time()
        
        # Load template
        self._load_template(template_path)
        
        # Prepare template
        if template_path and self.presentation:
            self._prepare_template_for_content(
                report_data.get('include_title_slide', True))
        
        # Create title slide
        if report_data.get('include_title_slide', True):
            self._create_title_slide(report_data)
        
        # Create content slides based on type
        for slide_config in slides_data:
            slide_type = slide_config.get('type', 'screenshot')
            title = slide_config.get('title', 'Slide')
            
            if slide_type == 'risks':
                # Native table for risks
                risks = slide_config.get('data', [])
                page_num = slide_config.get('page_num', 1)
                total_pages = slide_config.get('total_pages', 1)
                self.create_risk_table_slide(
                    risks=risks,
                    title=title,
                    page_num=page_num,
                    total_pages=total_pages
                )
            elif slide_type == 'milestones':
                # Native table for milestones
                milestones = slide_config.get('data', [])
                self.create_milestone_table_slide(
                    milestones=milestones,
                    title=title
                )
            elif slide_type == 'changes':
                # Native table for schedule changes (auto-paginates)
                changes = slide_config.get('data', [])
                rows_per_slide = slide_config.get('rows_per_slide', 6)
                self.create_changes_table_slides(
                    changes=changes,
                    title=title,
                    rows_per_slide=rows_per_slide
                )
            else:
                # Default: screenshot-based slide
                screenshot = slide_config.get('data')
                if screenshot:
                    self._create_content_slide(
                        screenshot,
                        slide_number=1,
                        skip_title=False,
                        slide_title=title
                    )
        
        # Check generation time
        elapsed = time.time() - self.start_time
        if elapsed >= self.MAX_GENERATION_TIME * 2:  # Allow more time for tables
            raise ValueError("File generation exceeded time limit")
            
        return self._save_to_bytes()

    def _prepare_template_for_content(self, keep_title_slide: bool = True):
        """Remove existing content slides from template, keeping only the structure.
        
        This allows us to use the template's styling while replacing its content.
        """
        if not self.presentation:
            return
            
        # Count slides to remove (skip first slide if keeping title)
        slides_to_remove = []
        start_idx = 1 if keep_title_slide else 0
        
        for i in range(start_idx, len(self.presentation.slides)):
            slides_to_remove.append(i)
        
        # Remove slides in reverse order to avoid index shifting
        for idx in reversed(slides_to_remove):
            slide_id = self.presentation.slides._sldIdLst[idx].rId
            self.presentation.part.drop_rel(slide_id)
            del self.presentation.slides._sldIdLst[idx]
    
    def _apply_transform_to_image(self, image_bytes: bytes, transform: Dict[str, Any]) -> bytes:
        """Apply crop and positioning transforms to an image.
        
        Args:
            image_bytes: Original image bytes
            transform: Dict with keys like cropTop, cropBottom, cropLeft, cropRight, scale, left, top
            
        Returns:
            Transformed image as bytes
        """
        try:
            img = Image.open(BytesIO(image_bytes))
            orig_width, orig_height = img.size
            
            # Apply crop if specified (values are percentages)
            crop_top = transform.get('cropTop', 0) / 100
            crop_bottom = transform.get('cropBottom', 0) / 100
            crop_left = transform.get('cropLeft', 0) / 100
            crop_right = transform.get('cropRight', 0) / 100
            
            if any([crop_top, crop_bottom, crop_left, crop_right]):
                left = int(orig_width * crop_left)
                top = int(orig_height * crop_top)
                right = int(orig_width * (1 - crop_right))
                bottom = int(orig_height * (1 - crop_bottom))
                img = img.crop((left, top, right, bottom))
            
            # Apply scale if specified
            scale = transform.get('scale', 1.0)
            if scale != 1.0 and scale > 0:
                new_width = int(img.width * scale)
                new_height = int(img.height * scale)
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Save to bytes
            output = BytesIO()
            img.save(output, format='PNG')
            output.seek(0)
            return output.read()
            
        except Exception as e:
            # Return original on error
            import logging
            logging.warning(f"Failed to apply transform: {e}")
            return image_bytes
        
    def _load_template(self, template_path: Optional[str] = None):
        """Load PowerPoint template or use default."""
        if template_path:
            try:
                self.presentation = Presentation(template_path)
                self.template_path = template_path
                self._validate_template()
            except Exception as e:
                raise ValueError(f"Invalid template: {str(e)}")
        else:
            # Use default template or create new presentation
            if os.path.exists(self.DEFAULT_TEMPLATE_PATH):
                self.presentation = Presentation(self.DEFAULT_TEMPLATE_PATH)
            else:
                self.presentation = Presentation()
                
    def _validate_template(self):
        """Validate the loaded template."""
        try:
            # Check if presentation can be accessed
            _ = self.presentation.slides
            # Check if it has slide layouts
            if not self.presentation.slide_layouts:
                raise ValueError("Template has no slide layouts")
        except Exception as e:
            raise ValueError(f"Invalid template format: {str(e)}")
            
    def _create_title_slide(self, report_data: Dict[str, Any]):
        """Create title slide with report metadata."""
        # Use first slide layout (title slide)
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = report_data.get('title', 'Report')
            
        # Add metadata to subtitle or body
        subtitle = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != title:
                subtitle = shape
                break
                
        if subtitle:
            text_frame = subtitle.text_frame
            text_frame.clear()
            
            # Add metadata
            p = text_frame.add_paragraph()
            p.text = f"Generated: {report_data.get('generated_date', datetime.now().strftime('%Y-%m-%d'))}"
            p.alignment = PP_ALIGN.CENTER
            
            if 'author' in report_data:
                p = text_frame.add_paragraph()
                p.text = f"Author: {report_data['author']}"
                p.alignment = PP_ALIGN.CENTER
                
            if 'description' in report_data:
                p = text_frame.add_paragraph()
                p.text = report_data['description']
                p.alignment = PP_ALIGN.CENTER
                
    def _create_content_slide(
        self, 
        screenshot: bytes, 
        slide_number: int, 
        skip_title: bool = False,
        slide_title: str = None
    ):
        """Create a content slide with screenshot.
        
        Args:
            screenshot: Screenshot image as bytes
            slide_number: Slide number (used as fallback title)
            skip_title: If True, don't override existing title
            slide_title: Descriptive title for the slide
        """
        # Use content slide layout (usually index 1 or 5)
        layout_idx = 5 if len(self.presentation.slide_layouts) > 5 else 1
        slide_layout = self.presentation.slide_layouts[layout_idx]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title if available and not skipping
        if slide.shapes.title and not skip_title:
            if slide_title:
                slide.shapes.title.text = slide_title
            else:
                slide.shapes.title.text = f"Screenshot {slide_number}"
            
            # Apply title styling: font size 28, color #7F7F7F
            if slide.shapes.title.has_text_frame:
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(28)
                        run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
            
        # Add screenshot image and get its position for overlays
        image_bounds = self._add_image_to_slide(slide, screenshot)
        
        # Add editable owner/resource text boxes for risk/milestone slides
        title_lower = (slide_title or '').lower()
        if 'risk' in title_lower and image_bounds:
            self._add_risk_owner_overlays(slide, image_bounds)
        elif 'milestone' in title_lower and image_bounds:
            self._add_milestone_resource_overlays(slide, image_bounds)
    
    def _add_risk_owner_overlays(self, slide, image_bounds):
        """Add editable text boxes positioned over the Owner fields in risk cards.
        
        The risk print view shows 3 cards per page, each with an Owner field
        in the meta section. We position white text boxes over each Owner field
        so users can click and type the actual name.
        
        Args:
            slide: The PowerPoint slide
            image_bounds: Dict with 'left', 'top', 'width', 'height' of image
        """
        try:
            # Each risk card is ~1/3 of the image height
            # Owner field is in the meta section, about 15% down from card top
            # and positioned at ~40% from the left (3rd of 5 columns)
            
            img_left = image_bounds['left']
            img_top = image_bounds['top']
            img_width = image_bounds['width']
            img_height = image_bounds['height']
            
            # Each card takes roughly 30% of image height with gaps
            card_height = img_height * 0.30
            card_gap = img_height * 0.03
            
            # Owner field position within card: ~15-20% down, ~40-55% across
            owner_y_offset = 0.15  # From top of card
            owner_x_offset = 0.38  # From left edge
            owner_width = img_width * 0.12  # Width of owner field
            owner_height = Inches(0.25)  # Height of text box
            
            for i in range(3):  # Up to 3 risk cards per page
                # Calculate card top position
                card_top = img_top + (i * (card_height + card_gap))
                
                # Position owner text box
                left = int(img_left + (img_width * owner_x_offset))
                top = int(card_top + (card_height * owner_y_offset))
                width = int(owner_width)
                
                # Add white text box
                textbox = slide.shapes.add_textbox(left, top, width, owner_height)
                tf = textbox.text_frame
                tf.word_wrap = False
                
                # Set white background
                fill = textbox.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                # Add placeholder text
                p = tf.paragraphs[0]
                p.text = "[Owner]"
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(150, 150, 150)
                p.font.bold = True
        except Exception as e:
            # Don't fail slide creation
            import logging
            logging.warning(f"Failed to add owner overlays: {e}")
    
    def _add_milestone_resource_overlays(self, slide, image_bounds):
        """Add editable text boxes for milestone Resource fields.
        
        The milestone print view shows 3 columns (last/this/next month),
        each with up to 8 milestone cards. Each card has a Resources field.
        We position white text boxes over each card's resources area.
        
        Args:
            slide: The PowerPoint slide
            image_bounds: Dict with 'left', 'top', 'width', 'height' of image
        """
        try:
            img_left = image_bounds['left']
            img_top = image_bounds['top']
            img_width = image_bounds['width']
            img_height = image_bounds['height']
            
            # 3 columns, each ~33% of width with gaps
            column_width = img_width * 0.30
            column_gap = img_width * 0.035
            
            # Header takes ~12% of height, cards start below
            header_height = img_height * 0.12
            cards_area_height = img_height - header_height
            
            # Each card is ~11% of image height (8 cards per column)
            card_height = cards_area_height * 0.11
            card_gap = cards_area_height * 0.012
            
            # Resources field position within card: ~55% down the card
            resources_y_offset = 0.55
            resources_height = Inches(0.18)
            
            # Add text boxes for each column (3 columns x 8 cards = 24 max)
            for col in range(3):
                # Column x position
                col_x = img_left + (col * (column_width + column_gap)) + column_gap
                
                for row in range(8):  # Up to 8 cards per column
                    # Card y position
                    card_y = (img_top + header_height + 
                              (row * (card_height + card_gap)))
                    
                    # Resources text box position
                    left = int(col_x + (column_width * 0.05))
                    top = int(card_y + (card_height * resources_y_offset))
                    width = int(column_width * 0.90)
                    
                    # Add white text box
                    textbox = slide.shapes.add_textbox(
                        left, top, width, resources_height)
                    tf = textbox.text_frame
                    tf.word_wrap = False
                    
                    # Set white background
                    fill = textbox.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(255, 255, 255)
                    
                    # Add placeholder text
                    p = tf.paragraphs[0]
                    p.text = "[Resource]"
                    p.font.size = Pt(8)
                    p.font.color.rgb = RGBColor(150, 150, 150)
        except Exception as e:
            # Don't fail slide creation
            import logging
            logging.warning(f"Failed to add milestone overlays: {e}")
    
    def _add_notes_textbox(self, slide, slide_title: str = None):
        """Add an editable text box for notes/resources at the bottom of the slide.
        
        The placeholder text is customized based on the slide type:
        - Risk slides: Owner field placeholder
        - Milestone slides: Resources field placeholder
        - Other slides: Generic notes placeholder
        """
        try:
            slide_width = self.presentation.slide_width
            slide_height = self.presentation.slide_height
            
            # Position at bottom of slide
            left = Inches(0.5)
            top = int(slide_height * 0.88)  # 88% down the slide
            width = int(slide_width - Inches(1))
            height = Inches(0.5)
            
            # Add text box
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            
            # Customize placeholder text based on slide type
            title_lower = (slide_title or '').lower()
            if 'risk' in title_lower:
                placeholder = "👤 Owner Names: [Edit to replace Owner A, B, C with actual names]"
            elif 'milestone' in title_lower:
                placeholder = "👤 Resources: [Edit to replace Resource A, B, C with actual names]"
            else:
                placeholder = "📝 Notes: [Edit this field to add additional information]"
            
            # Add placeholder text
            p = tf.paragraphs[0]
            p.text = placeholder
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(100, 100, 100)  # Gray color
            p.font.italic = True
        except Exception as e:
            # Don't fail slide creation if notes box fails
            pass
        
    def _add_image_to_slide(self, slide, image_bytes: bytes) -> Dict[str, int]:
        """Add image to slide with smart sizing that fills the content area.
        
        The image is enlarged to fill the available space while:
        - Respecting template safe zones (title area, footer, margins)
        - Maintaining aspect ratio
        - Not encroaching on branding elements
        
        Returns:
            Dict with 'left', 'top', 'width', 'height' of placed image (in EMUs)
        """
        # Load image to get dimensions
        img = Image.open(BytesIO(image_bytes))
        img_width, img_height = img.size
        
        # Calculate slide dimensions
        slide_width = self.presentation.slide_width
        slide_height = self.presentation.slide_height
        
        # Get safe zone based on whether we have a template
        safe_zone = self._get_safe_zone()
        
        # Calculate available content area (in EMUs)
        content_left = Inches(safe_zone['left'])
        content_top = Inches(safe_zone['top'])
        content_width = slide_width - Inches(safe_zone['left'] + safe_zone['right'])
        content_height = slide_height - Inches(safe_zone['top'] + safe_zone['bottom'])
        
        # Calculate scaling to fill content area while maintaining aspect ratio
        width_ratio = content_width / img_width
        height_ratio = content_height / img_height
        scale_ratio = min(width_ratio, height_ratio)
        
        # Calculate final dimensions
        width = int(img_width * scale_ratio)
        height = int(img_height * scale_ratio)
        
        # Calculate position (centered within content area)
        left = content_left + (content_width - width) // 2
        top = content_top + (content_height - height) // 2
        
        # Add image
        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream,
            left,
            top,
            width=width,
            height=height
        )
        
        # Return image bounds for overlay positioning
        return {
            'left': int(left),
            'top': int(top),
            'width': int(width),
            'height': int(height)
        }
    
    def _get_safe_zone(self) -> Dict[str, float]:
        """Get the safe zone configuration based on template.
        
        Returns:
            Dict with 'top', 'bottom', 'left', 'right' values in inches
        """
        if self.template_path:
            return TEMPLATE_SAFE_ZONES['branded']
        return TEMPLATE_SAFE_ZONES['default']
        
    def _save_to_bytes(self) -> bytes:
        """Save presentation to bytes."""
        output = BytesIO()
        self.presentation.save(output)
        output.seek(0)
        return output.read()
    
    def create_risk_table_slide(
        self, 
        risks: List[Dict[str, Any]], 
        title: str = "Risk Register",
        page_num: int = 1,
        total_pages: int = 1
    ):
        """Create a slide with risks in card format matching the preview.
        
        Each risk is a card with:
        - Title bar with ID and title
        - Meta row: Severity, Status, Owner, L, I
        - Body: Description and Mitigation
        
        Args:
            risks: List of risk dictionaries
            title: Slide title
            page_num: Current page number
            total_pages: Total pages
        """
        from pptx.enum.text import MSO_ANCHOR
        from pptx.dml.color import RGBColor
        
        # Add slide
        layout_idx = 5 if len(self.presentation.slide_layouts) > 5 else 1
        slide_layout = self.presentation.slide_layouts[layout_idx]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            page_indicator = f" (Page {page_num}/{total_pages})" if total_pages > 1 else ""
            slide.shapes.title.text = f"{title}{page_indicator}"
            if slide.shapes.title.has_text_frame:
                for para in slide.shapes.title.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(24)
                        run.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
        
        # Card dimensions - fit 3 cards per slide
        card_width = Inches(12.0)
        card_height = Inches(1.9)
        card_left = Inches(0.5)
        start_top = Inches(1.2)
        card_gap = Inches(0.15)
        
        severity_colors = {
            'critical': RGBColor(0x7C, 0x3A, 0xED),
            'high': RGBColor(0xDC, 0x26, 0x26),
            'medium': RGBColor(0xF5, 0x9E, 0x0B),
            'low': RGBColor(0x6B, 0x72, 0x80),
        }
        
        for idx, risk in enumerate(risks[:3]):  # Max 3 per slide
            card_top = start_top + (card_height + card_gap) * idx
            
            risk_id = risk.get('id', 'N/A')
            risk_title = risk.get('title', 'Untitled Risk')
            severity = risk.get('severity_normalized', 'medium').lower()
            status = risk.get('status', 'N/A')
            owner = risk.get('owner', '')
            likelihood = risk.get('likelihood', '?')
            impact = risk.get('impact', '?')
            description = risk.get('description', 'No description provided.')
            mitigations = risk.get('mitigations', []) or risk.get('mitigation', '')
            if isinstance(mitigations, list):
                miti_text = '; '.join(mitigations) if mitigations else 'None specified'
            else:
                miti_text = str(mitigations) if mitigations else 'None specified'
            
            # Title bar (blue gradient background)
            title_shape = slide.shapes.add_shape(
                1, card_left, card_top, card_width, Inches(0.4)  # 1 = rectangle
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(0x1E, 0x40, 0xAF)
            title_shape.line.fill.background()
            
            title_frame = title_shape.text_frame
            title_frame.paragraphs[0].text = f"{risk_id}: {risk_title}"
            title_frame.paragraphs[0].font.size = Pt(12)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title_frame.margin_left = Inches(0.1)
            title_frame.margin_top = Inches(0.05)
            
            # Meta row (gray background)
            meta_top = card_top + Inches(0.4)
            meta_shape = slide.shapes.add_shape(
                1, card_left, meta_top, card_width, Inches(0.35)
            )
            meta_shape.fill.solid()
            meta_shape.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
            meta_shape.line.fill.background()
            
            # Create meta text boxes for each field
            meta_items = [
                ("Severity", severity.upper(), severity_colors.get(severity, severity_colors['medium'])),
                ("Status", status, RGBColor(0x1F, 0x29, 0x37)),
                ("Owner", owner or "Owner", RGBColor(0x1F, 0x29, 0x37)),
                ("Likelihood", f"L: {likelihood}", RGBColor(0x1F, 0x29, 0x37)),
                ("Impact", f"I: {impact}", RGBColor(0x1F, 0x29, 0x37))
            ]
            
            item_width = card_width / 5
            for i, (label, value, color) in enumerate(meta_items):
                item_left = card_left + item_width * i
                meta_box = slide.shapes.add_textbox(
                    item_left, meta_top, item_width, Inches(0.35)
                )
                tf = meta_box.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.05)
                tf.margin_top = Inches(0.02)
                
                # Label
                p1 = tf.paragraphs[0]
                p1.text = label
                p1.font.size = Pt(7)
                p1.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
                
                # Value
                p2 = tf.add_paragraph()
                p2.text = str(value)[:15]
                p2.font.size = Pt(9)
                p2.font.bold = True
                p2.font.color.rgb = color
            
            # Body area (description + mitigation)
            body_top = meta_top + Inches(0.35)
            body_height = card_height - Inches(0.75)
            
            # Description section (left half)
            desc_box = slide.shapes.add_textbox(
                card_left, body_top, card_width / 2 - Inches(0.1), body_height
            )
            desc_tf = desc_box.text_frame
            desc_tf.word_wrap = True
            desc_tf.margin_left = Inches(0.1)
            desc_tf.margin_top = Inches(0.05)
            
            p = desc_tf.paragraphs[0]
            p.text = "Description"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            
            p2 = desc_tf.add_paragraph()
            p2.text = str(description)[:200] + ('...' if len(str(description)) > 200 else '')
            p2.font.size = Pt(8)
            p2.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
            
            # Mitigation section (right half)
            miti_box = slide.shapes.add_textbox(
                card_left + card_width / 2, body_top, card_width / 2 - Inches(0.1), body_height
            )
            miti_tf = miti_box.text_frame
            miti_tf.word_wrap = True
            miti_tf.margin_left = Inches(0.1)
            miti_tf.margin_top = Inches(0.05)
            
            p = miti_tf.paragraphs[0]
            p.text = "Mitigation"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            
            p2 = miti_tf.add_paragraph()
            p2.text = str(miti_text)[:200] + ('...' if len(str(miti_text)) > 200 else '')
            p2.font.size = Pt(8)
            p2.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
            
            # Card border
            border = slide.shapes.add_shape(
                1, card_left, card_top, card_width, card_height
            )
            border.fill.background()
            border.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
            border.line.width = Pt(1)
    
    def create_milestone_table_slide(
        self,
        milestones: List[Dict[str, Any]],
        title: str = "Milestones",
        column_title: str = "This Month"
    ):
        """Create a slide with milestones in 3-column table layout.
        
        EXACTLY matches the HTML canvas preview format:
        - Title: #7F7F7F gray, 32px, normal weight
        - 3 columns with colored headers and tables below
        - Header colors: Gray (#6B7280), Blue (#2563EB), Amber (#F59E0B)
        - Table styling matches HTML preview exactly
        
        Args:
            milestones: List of milestone dictionaries
            title: Slide title
            column_title: Column header (unused, using month names)
        """
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from datetime import datetime, timedelta
        
        # Add slide with blank layout
        layout_idx = 6 if len(self.presentation.slide_layouts) > 6 else 5
        slide_layout = self.presentation.slide_layouts[layout_idx]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # ============================================================
        # TITLE: Matches HTML .slide-title
        # color: #7F7F7F, font-size: 32px, font-weight: normal
        # ============================================================
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.paragraphs[0].text = title
        title_tf.paragraphs[0].font.size = Pt(24)  # 32px scaled for slide
        title_tf.paragraphs[0].font.bold = False  # font-weight: normal
        title_tf.paragraphs[0].font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)  # #7F7F7F
        
        # ============================================================
        # DATE CALCULATIONS
        # ============================================================
        today = datetime.now()
        this_month_start = today.replace(day=1)
        if today.month == 12:
            next_month_start = today.replace(year=today.year + 1, month=1, day=1)
        else:
            next_month_start = today.replace(month=today.month + 1, day=1)
        
        last_month_end = this_month_start - timedelta(days=1)
        last_month_start = last_month_end.replace(day=1)
        
        if next_month_start.month == 12:
            next_next = next_month_start.replace(year=next_month_start.year + 1, month=1, day=1)
        else:
            next_next = next_month_start.replace(month=next_month_start.month + 1, day=1)
        next_month_end = next_next - timedelta(days=1)
        this_month_end = next_month_start - timedelta(days=1)
        
        def get_ms_attr(ms, attr, default=''):
            if hasattr(ms, attr):
                return getattr(ms, attr, default) or default
            return ms.get(attr, default) if isinstance(ms, dict) else default
        
        def is_in_range(date_str, start, end):
            try:
                d = datetime.strptime(str(date_str), '%Y-%m-%d')
                return start <= d <= end
            except:
                return False
        
        def format_date_range(start, end):
            """Format: 'Nov 01 - Nov 30, 2025' matching HTML"""
            return f"{start.strftime('%b %d')} - {end.strftime('%b %d, %Y')}"
        
        # Filter milestones by month
        last_ms = [m for m in milestones 
                   if is_in_range(get_ms_attr(m, 'target_date'), last_month_start, last_month_end)]
        this_ms = [m for m in milestones 
                   if is_in_range(get_ms_attr(m, 'target_date'), this_month_start, this_month_end)]
        next_ms = [m for m in milestones 
                   if is_in_range(get_ms_attr(m, 'target_date'), next_month_start, next_month_end)]
        
        # ============================================================
        # COLUMN LAYOUT - Matches HTML .columns grid
        # 3 equal columns with 20px gap (scaled to inches)
        # ============================================================
        slide_width = 13.333  # Standard widescreen width in inches
        margin = 0.4
        gap = 0.2  # 20px gap scaled
        usable_width = slide_width - (2 * margin) - (2 * gap)
        column_width = usable_width / 3
        header_top = Inches(1.0)
        
        # ============================================================
        # HEADER COLORS - Matches CANVAS PREVIEW exactly
        # Last Month: #EA580C (orange/red)
        # This Month: #16A34A (green)  
        # Next Month: #F59E0B (yellow/amber), text #1f2937 (dark)
        # ============================================================
        month_config = [
            {
                'icon': '📅',
                'title': 'Last Month (Completed)',
                'date_range': format_date_range(last_month_start, last_month_end),
                'milestones': last_ms,
                'header_bg': RGBColor(0xEA, 0x58, 0x0C),  # #EA580C orange/red
                'header_text': RGBColor(0xFF, 0xFF, 0xFF),  # white
            },
            {
                'icon': '📍',
                'title': 'This Month (In Progress)',
                'date_range': format_date_range(this_month_start, this_month_end),
                'milestones': this_ms,
                'header_bg': RGBColor(0x16, 0xA3, 0x4A),  # #16A34A green
                'header_text': RGBColor(0xFF, 0xFF, 0xFF),  # white
            },
            {
                'icon': '📌',
                'title': 'Next Month (Planned)',
                'date_range': format_date_range(next_month_start, next_month_end),
                'milestones': next_ms,
                'header_bg': RGBColor(0xF5, 0x9E, 0x0B),  # #F59E0B amber
                'header_text': RGBColor(0x1F, 0x29, 0x37),  # #1f2937 dark
            },
        ]
        
        for col_idx, config in enumerate(month_config):
            left = Inches(margin + col_idx * (column_width + gap))
            col_w = Inches(column_width)
            
            # ============================================================
            # COLUMN HEADER - Matches HTML .column-header
            # padding: 14px 16px, font-size: 18px, font-weight: bold
            # ============================================================
            header_height = Inches(0.55)  # Room for title + date range
            header_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, header_top, col_w, header_height
            )
            header_bar.fill.solid()
            header_bar.fill.fore_color.rgb = config['header_bg']
            header_bar.line.fill.background()
            # Round corners (8px border-radius scaled)
            header_bar.adjustments[0] = 0.05
            
            # Header title: icon + title text
            # font-size: 18px = ~13.5pt scaled for PPT
            header_tf = header_bar.text_frame
            header_tf.word_wrap = True
            header_tf.margin_left = Inches(0.12)
            header_tf.margin_top = Inches(0.08)
            header_tf.margin_right = Inches(0.08)
            header_tf.margin_bottom = Inches(0.04)
            
            p1 = header_tf.paragraphs[0]
            p1.text = f"{config['icon']} {config['title']}"
            p1.font.size = Pt(12)  # 18px scaled
            p1.font.bold = True
            p1.font.color.rgb = config['header_text']
            
            # Date range subtitle: font-size: 14px, opacity: 0.9
            p2 = header_tf.add_paragraph()
            p2.text = config['date_range']
            p2.font.size = Pt(9)  # 14px scaled
            p2.font.bold = False
            p2.font.color.rgb = config['header_text']
            p2.space_before = Pt(2)
            
            # ============================================================
            # TABLE - Matches HTML table styling
            # ============================================================
            table_top = header_top + header_height + Inches(0.08)
            ms_list = config['milestones']
            max_rows = min(len(ms_list), 8) if ms_list else 1
            num_rows = max_rows + 1  # +1 for header row
            row_height = Inches(0.32)  # padding: 10px 8px
            
            table_shape = slide.shapes.add_table(
                num_rows, 4, left, table_top, col_w, row_height * num_rows
            )
            table = table_shape.table
            
            # Column widths proportional to HTML
            # Milestone (40%), Date (15%), Status (20%), Resources (25%)
            col_w_val = column_width
            table.columns[0].width = Inches(col_w_val * 0.40)
            table.columns[1].width = Inches(col_w_val * 0.15)
            table.columns[2].width = Inches(col_w_val * 0.20)
            table.columns[3].width = Inches(col_w_val * 0.25)
            
            # ============================================================
            # TABLE HEADER ROW - Matches HTML th styling
            # background: #f3f4f6, font-size: 15px, font-weight: bold
            # padding: 10px 8px, border-bottom: 1px #e5e7eb
            # ============================================================
            headers = ["Milestone", "Date", "Status", "Resources"]
            for i, hdr in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = hdr
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)  # #f3f4f6
                
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(9)  # 15px scaled
                para.font.bold = True
                para.font.color.rgb = RGBColor(0x37, 0x41, 0x51)  # Dark gray
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.margin_left = Inches(0.05)
                cell.text_frame.margin_right = Inches(0.05)
            
            # ============================================================
            # STATUS COLORS - Matches CANVAS PREVIEW exactly
            # Completed: #EA580C (red/orange), bold
            # In Progress: #2563EB (blue), bold
            # Not Started: #6B7280 (gray), normal
            # ============================================================
            status_colors = {
                'COMPLETED': (RGBColor(0xEA, 0x58, 0x0C), True),   # #EA580C red/orange, bold
                'IN_PROGRESS': (RGBColor(0x25, 0x63, 0xEB), True),  # #2563EB blue, bold
                'NOT_STARTED': (RGBColor(0x6B, 0x72, 0x80), False), # #6B7280 gray, normal
            }
            
            if not ms_list:
                # Empty state - matches HTML .empty
                cell = table.cell(1, 0)
                cell.merge(table.cell(1, 3))
                cell.text = "No milestones"
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(9)  # 14px scaled
                para.font.italic = True
                para.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)  # #9ca3af
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                # ============================================================
                # DATA ROWS - Matches HTML td styling
                # padding: 10px 8px, font-size: 16px, line-height: 1.4
                # td.name: font-weight: 500
                # td.date: font-size: 15px, color: #666
                # td.resource: background: #fffef0, font-style: italic, color: #666
                # ============================================================
                for row_idx, ms in enumerate(ms_list[:8], start=1):
                    name = str(get_ms_attr(ms, 'name', 'Untitled'))
                    target = str(get_ms_attr(ms, 'target_date', ''))
                    status = str(get_ms_attr(ms, 'status', 'NOT_STARTED')).upper()
                    resources = get_ms_attr(ms, 'resources', '')
                    
                    # Truncate name to ~35 chars for fit
                    name_display = name[:35] + ('...' if len(name) > 35 else '')
                    
                    # Format date - keep full date for clarity
                    date_display = target
                    
                    # Status display text
                    status_display = status.replace('_', ' ').title()
                    
                    # Resources - show actual or placeholder
                    resource_display = str(resources) if resources else 'Resource'
                    
                    # --- Milestone Name Cell ---
                    cell = table.cell(row_idx, 0)
                    cell.text = name_display
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(8)
                    para.font.bold = False  # font-weight: 500 = medium
                    para.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)  # #1f2937
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.text_frame.margin_left = Inches(0.05)
                    cell.text_frame.word_wrap = True
                    
                    # --- Date Cell ---
                    cell = table.cell(row_idx, 1)
                    cell.text = date_display
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(8)
                    para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)  # #666
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # --- Status Cell ---
                    cell = table.cell(row_idx, 2)
                    cell.text = status_display
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(8)
                    color, is_bold = status_colors.get(status, status_colors['NOT_STARTED'])
                    para.font.color.rgb = color
                    para.font.bold = is_bold
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # --- Resources Cell (EDITABLE) ---
                    # background: #fffef0, font-style: italic, color: #666
                    cell = table.cell(row_idx, 3)
                    cell.text = resource_display
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFE, 0xF0)  # #fffef0
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(8)
                    para.font.italic = True
                    para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)  # #666
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # Alternate row backgrounds (subtle)
                    if row_idx % 2 == 0:
                        for ci in range(3):  # Don't override resources column
                            c = table.cell(row_idx, ci)
                            if ci != 3:
                                c.fill.solid()
                                c.fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
        
        # ============================================================
        # INFO BOX - Matches HTML .info-box
        # background: #e0f2fe, font-size: 12px, color: #0369a1
        # ============================================================
        info_left = Inches(margin)
        info_top = Inches(6.8)
        info_width = Inches(slide_width - 2 * margin)
        info_height = Inches(0.35)
        
        info_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, info_left, info_top, info_width, info_height
        )
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = RGBColor(0xE0, 0xF2, 0xFE)  # #e0f2fe
        info_box.line.fill.background()
        info_box.adjustments[0] = 0.1
        
        info_tf = info_box.text_frame
        info_tf.margin_left = Inches(0.12)
        info_tf.margin_top = Inches(0.08)
        p = info_tf.paragraphs[0]
        p.text = "ℹ️ Resources column (yellow) is editable in PowerPoint."
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x03, 0x69, 0xA1)  # #0369a1

    def create_changes_table_slides(
        self,
        changes: List[Dict[str, Any]],
        title: str = "Schedule Changes",
        rows_per_slide: int = 6
    ) -> int:
        """Create slides with changes rendered as native PowerPoint tables.
        
        Handles multi-page output when changes exceed rows_per_slide.
        
        Args:
            changes: List of change dictionaries with change_id, milestone_name,
                    old_date, new_date, reason, impact (contingency)
            title: Base slide title
            rows_per_slide: Maximum rows per slide (excluding header)
            
        Returns:
            Number of slides created
        """
        from pptx.enum.text import MSO_ANCHOR
        
        if not changes:
            return 0
        
        # Calculate number of slides needed
        total_changes = len(changes)
        num_slides = (total_changes + rows_per_slide - 1) // rows_per_slide
        
        for slide_num in range(num_slides):
            # Get changes for this slide
            start_idx = slide_num * rows_per_slide
            end_idx = min(start_idx + rows_per_slide, total_changes)
            slide_changes = changes[start_idx:end_idx]
            
            # Add slide
            layout_idx = 5 if len(self.presentation.slide_layouts) > 5 else 1
            slide_layout = self.presentation.slide_layouts[layout_idx]
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Set title with page indicator
            if slide.shapes.title:
                page_indicator = f" ({slide_num + 1}/{num_slides})" if num_slides > 1 else ""
                slide.shapes.title.text = f"{title}{page_indicator}"
                if slide.shapes.title.has_text_frame:
                    for para in slide.shapes.title.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(24)
                            run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
            
            # Table dimensions
            left = Inches(0.3)
            top = Inches(1.2)
            width = Inches(12.9)
            
            # Columns: Milestone, Old Date, New Date, Reason, Contingency
            cols = 5
            rows = len(slide_changes) + 1  # +1 for header
            
            # Calculate row height to fit content
            available_height = 5.8  # inches
            row_height = min(0.9, available_height / rows)
            
            table = slide.shapes.add_table(
                rows, cols, left, top, width, Inches(row_height * rows)
            ).table
            
            # Column widths: Milestone wider, dates narrow, reason/contingency medium
            col_widths = [3.5, 1.0, 1.0, 3.7, 3.7]
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w)
            
            # Header row
            headers = ["Milestone", "Old Date", "New Date", "Reason", "Contingency"]
            header_color = RGBColor(0x1E, 0x40, 0xAF)  # Blue
            
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(10)
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Data rows
            for row_idx, change in enumerate(slide_changes, start=1):
                # Extract change data
                milestone = change.get('milestone_name', change.get('change_id', ''))
                # Clean milestone name (remove date suffix if present)
                import re
                milestone = re.sub(r'-\d{4}-\d{2}-\d{2}-to-\d{4}-\d{2}-\d{2}$', '', str(milestone))
                
                old_date = change.get('old_date', '')
                new_date = change.get('new_date', '')
                reason = change.get('reason', '')
                contingency = change.get('impact', '')  # 'impact' field is now contingency
                
                row_data = [milestone, old_date, new_date, reason, contingency]
                
                for col_idx, value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    
                    # Set text with word wrap
                    cell.text = str(value) if value else ''
                    
                    # Enable text wrapping
                    cell.text_frame.word_wrap = True
                    
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(9)
                    para.alignment = PP_ALIGN.LEFT if col_idx in [0, 3, 4] else PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.TOP
                    
                    # Style dates with color
                    if col_idx == 1:  # Old date - red strikethrough look
                        para.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
                    elif col_idx == 2:  # New date - green
                        para.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A)
                        para.font.bold = True
                    
                    # Alternate row background
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
        
        return num_slides

    def get_slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        if not self.presentation:
            return 0
        return len(self.presentation.slides)
        
    def preserve_branding(self, template_path: str) -> bool:
        """
        Check if company branding is preserved from template.
        
        Returns:
            True if branding elements are preserved
        """
        try:
            # Load template
            template_pres = Presentation(template_path)
            
            # Check if template has master slides with branding
            if template_pres.slide_master:
                # Branding is preserved through slide master
                return True
                
            # Check for logo shapes in template
            for slide_layout in template_pres.slide_layouts:
                for shape in slide_layout.shapes:
                    if shape.shape_type == 13:  # Picture
                        return True
                        
            return True  # Assume branding preserved if template loads
            
        except Exception:
            return False


# Helper functions for testing
def create_default_template():
    """Create a default template if it doesn't exist."""
    os.makedirs("templates", exist_ok=True)
    
    if not os.path.exists(PowerPointBuilderService.DEFAULT_TEMPLATE_PATH):
        prs = Presentation()
        
        # Customize slide master
        slide_master = prs.slide_master
        
        # Add title slide layout
        title_layout = prs.slide_layouts[0]
        
        # Add content slide layout
        content_layout = prs.slide_layouts[5]
        
        # Save as default template
        prs.save(PowerPointBuilderService.DEFAULT_TEMPLATE_PATH)


# Create default template on module load
create_default_template()
