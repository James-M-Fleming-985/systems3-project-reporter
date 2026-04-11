"""
PowerPoint Export Service - Generate PPTX presentations from project data
"""
from pathlib import Path
from datetime import datetime, timedelta
from typing import List, Dict, Any
from io import BytesIO
from calendar import monthrange
import logging

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None

from models import Project, Milestone


class PowerPointExporter:
    """Export project data to PowerPoint presentation with enhanced layout"""
    
    def __init__(self, template_path: Path = None):
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PowerPoint export. "
                "Install with: pip install python-pptx"
            )
        self.template_path = template_path
    
    def create_presentation(self, projects: List[Project]) -> BytesIO:
        """
        Create enhanced PowerPoint presentation:
        1. Title Page
        2. Roadmap/Gantt Page (all projects timeline)
        3. Milestones Page (top 5-6 milestones in app format)
        4. Risks Page (app format)
        5. Change Management Page (app format)
        
        Args:
            projects: List of Project objects to include
            
        Returns:
            BytesIO buffer containing the PPTX file
        """
        # Use custom template if provided, otherwise create blank
        if self.template_path and self.template_path.exists():
            logger.info(f"Using custom template: {self.template_path}")
            prs = Presentation(str(self.template_path))
        else:
            logger.info("Using default blank presentation")
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        
        # Collect all milestones, risks, and changes from all projects
        all_milestones = []
        all_risks = []
        all_changes = []
        
        for project in projects:
            for milestone in project.milestones:
                try:
                    milestone_data = {
                        'name': milestone.name or 'Unnamed Milestone',
                        'project': project.project_name or 'Unknown Project',
                        'parent_project': milestone.parent_project,
                        'target_date': milestone.target_date or '2025-01-01',
                        'status': milestone.status or 'NOT_STARTED',
                        'completion_percentage': milestone.completion_percentage,
                        'resources': milestone.resources,
                        'notes': milestone.notes
                    }
                    all_milestones.append(milestone_data)
                except AttributeError:
                    continue
            
            for risk in project.risks:
                try:
                    all_risks.append({
                        'project': project.project_name or 'Unknown Project',
                        'description': risk.description or 'No description',
                        'severity': risk.severity or 'LOW',
                        'status': risk.status or 'OPEN',
                        'mitigation': risk.mitigation or 'No mitigation plan'
                    })
                except AttributeError:
                    continue
            
            for change in project.changes:
                try:
                    all_changes.append({
                        'project': project.project_name or 'Unknown Project',
                        'date': change.date or '2025-01-01',
                        'old_date': change.old_date or '2025-01-01',
                        'new_date': change.new_date or '2025-01-01',
                        'reason': change.reason or 'No reason provided',
                        'impact': change.impact
                    })
                except AttributeError:
                    continue
        
        # Sort milestones by date (ascending) - matching app default
        # Only sort if we have valid dates
        if all_milestones:
            try:
                all_milestones.sort(key=lambda m: datetime.strptime(m['target_date'], '%Y-%m-%d'))
            except (ValueError, KeyError, TypeError):
                # If sorting fails, keep original order
                pass
        
        # 1. Title slide
        self._add_title_slide(prs, projects, all_milestones)
        
        # 2. Gantt/Roadmap page (visual snapshot of all projects)
        self._add_gantt_roadmap_slide(prs, projects)
        
        logger.info(f"Collected {len(all_milestones)} milestones, {len(all_risks)} risks, {len(all_changes)} changes")
        
        # 3. Milestones page (3-column time-based view)
        self._add_milestones_slide_app_format(prs, all_milestones)  # Pass all, method filters by month
        
        # 4. Risks page (app format)
        self._add_risks_slide_app_format(prs, all_risks)  # Always create, even if empty
        
        # 5. Change Management page (app format)
        if all_changes:
            self._add_changes_slide_app_format(prs, all_changes)
        
        # 6. Financial slides (if financial data exists)
        self._add_financial_slides(prs)
        
        # Save to BytesIO
        buffer = BytesIO()
        try:
            prs.save(buffer)
            buffer.seek(0)
            logger.info(f"PowerPoint created successfully: {buffer.getbuffer().nbytes} bytes")
            return buffer
        except Exception as e:
            logger.error(f"Failed to save PowerPoint: {e}")
            raise
    
    def _add_title_slide(self, prs, projects, all_milestones):
        """Add enhanced title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        completed = sum(1 for m in all_milestones if m['status'] == 'COMPLETED')
        
        title.text = "Systems³ Project Reporter"
        subtitle.text = (
            f"Project Status Report\n"
            f"Generated: {datetime.now().strftime('%B %d, %Y')}\n"
            f"{len(projects)} Active Projects | "
            f"{len(all_milestones)} Milestones ({completed} Completed)"
        )
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
        
        # Style subtitle
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
    
    def _add_gantt_roadmap_slide(self, prs, projects):
        """Add Gantt roadmap - clean table with project bars"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "📊 Program Roadmap - Project Overview"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
        
        if not projects:
            return
        
        # Filter projects with valid dates
        valid_projects = []
        for project in projects:
            try:
                start = datetime.strptime(str(project.start_date), '%Y-%m-%d')
                end = datetime.strptime(str(project.target_completion), '%Y-%m-%d')
                valid_projects.append({
                    'name': str(getattr(project, 'project_name', 'Unknown'))[:40],
                    'start': start,
                    'end': end,
                    'completion': int(getattr(project, 'completion_percentage', 0) or 0)
                })
            except:
                continue
        
        if not valid_projects:
            return
        
        # Create table: Project names + timeline visualization
        rows = len(valid_projects) + 1
        cols = 2  # Project name | Timeline bar
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.5)
        ).table
        
        # Set column widths
        table.columns[0].width = Inches(3.5)  # Project names
        table.columns[1].width = Inches(5.5)  # Timeline
        
        # Header row
        header_cells = ['Project', 'Timeline & Completion']
        for col_idx, header in enumerate(header_cells):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(249, 250, 251)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
        
        # Data rows
        for idx, proj in enumerate(valid_projects, 1):
            # Project name
            name_cell = table.cell(idx, 0)
            name_cell.text = proj['name']
            name_cell.text_frame.paragraphs[0].font.size = Pt(9)
            name_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            name_cell.vertical_anchor = 1  # Middle
            
            # Timeline description (since visual bars are complex in tables)
            timeline_cell = table.cell(idx, 1)
            start_str = proj['start'].strftime('%b %Y')
            end_str = proj['end'].strftime('%b %Y')
            duration_days = (proj['end'] - proj['start']).days
            duration_months = max(1, duration_days // 30)
            
            timeline_text = (
                f"{start_str} → {end_str}  |  "
                f"{duration_months} months  |  "
                f"{proj['completion']}% Complete"
            )
            
            p = timeline_cell.text_frame.paragraphs[0]
            p.text = timeline_text
            p.font.size = Pt(9)
            
            # Color code by completion
            if proj['completion'] >= 100:
                p.font.color.rgb = RGBColor(34, 197, 94)  # Green
                p.font.bold = True
            elif proj['completion'] >= 50:
                p.font.color.rgb = RGBColor(59, 130, 246)  # Blue
            else:
                p.font.color.rgb = RGBColor(107, 114, 128)  # Gray
            
            timeline_cell.vertical_anchor = 1  # Middle
            
            # Alternating row colors
            if idx % 2 == 0:
                for col in range(cols):
                    table.cell(idx, col).fill.solid()
                    table.cell(idx, col).fill.fore_color.rgb = RGBColor(249, 250, 251)
    
    def _add_milestones_slide_app_format(self, prs, milestones):
        """Add milestones slide - clean table format organized by time period"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "📍 Upcoming Milestones"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
        
        # Calculate date ranges
        today = datetime.now()
        this_month_start = datetime(today.year, today.month, 1)
        
        # Categorize milestones
        last_month = []
        this_month = []
        next_month = []
        
        for milestone in milestones:
            try:
                target = datetime.strptime(str(milestone.get('target_date', '')), '%Y-%m-%d')
                if target < this_month_start:
                    last_month.append(milestone)
                elif target.month == today.month and target.year == today.year:
                    this_month.append(milestone)
                else:
                    next_month.append(milestone)
            except:
                this_month.append(milestone)
        
        # Create comprehensive table
        all_display_milestones = (
            [(m, 'Last Month') for m in last_month[:5]] +
            [(m, 'This Month') for m in this_month[:5]] +
            [(m, 'Next Month') for m in next_month[:5]]
        )
        
        if not all_display_milestones:
            no_data = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            no_data.text_frame.text = "No milestones"
            no_data.text_frame.paragraphs[0].font.size = Pt(18)
            no_data.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            no_data.text_frame.paragraphs[0].font.color.rgb = RGBColor(107, 114, 128)
            return
        
        # Create table
        rows = len(all_display_milestones) + 1
        cols = 5  # Period, Milestone, Project, Target Date, Status
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.0),
            Inches(9), Inches(6)
        ).table
        
        # Set column widths
        table.columns[0].width = Inches(1.2)  # Period
        table.columns[1].width = Inches(3.0)  # Milestone
        table.columns[2].width = Inches(2.2)  # Project
        table.columns[3].width = Inches(1.4)  # Date
        table.columns[4].width = Inches(1.2)  # Status
        
        # Header row
        headers = ['Period', 'Milestone', 'Project', 'Target Date', 'Status']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(249, 250, 251)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
        
        # Data rows
        for idx, (milestone, period) in enumerate(all_display_milestones, 1):
            # Period
            period_cell = table.cell(idx, 0)
            if period == 'This Month':
                period_cell.text = "📍 This"
            elif period == 'Last Month':
                period_cell.text = "📅 Last"
            else:
                period_cell.text = "🔜 Next"
            period_cell.text_frame.paragraphs[0].font.size = Pt(9)
            period_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            period_cell.vertical_anchor = 1
            
            # Milestone name
            name_cell = table.cell(idx, 1)
            name_cell.text = str(milestone.get('name', 'Unnamed'))[:50]
            name_cell.text_frame.paragraphs[0].font.size = Pt(9)
            name_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            name_cell.text_frame.word_wrap = True
            name_cell.vertical_anchor = 1
            
            # Project
            proj_cell = table.cell(idx, 2)
            project = milestone.get('parent_project') or milestone.get('project', '')
            proj_cell.text = str(project)[:35]
            proj_cell.text_frame.paragraphs[0].font.size = Pt(9)
            proj_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            proj_cell.text_frame.word_wrap = True
            proj_cell.vertical_anchor = 1
            
            # Target date
            date_cell = table.cell(idx, 3)
            try:
                target_date = datetime.strptime(str(milestone.get('target_date', '')), '%Y-%m-%d')
                date_cell.text = target_date.strftime('%b %d, %Y')
            except:
                date_cell.text = str(milestone.get('target_date', ''))[:12]
            date_cell.text_frame.paragraphs[0].font.size = Pt(9)
            date_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            date_cell.vertical_anchor = 1
            
            # Status
            status_cell = table.cell(idx, 4)
            status = milestone.get('status', 'NOT_STARTED')
            
            if status == 'COMPLETED':
                status_cell.text = "✓ Done"
                status_cell.fill.solid()
                status_cell.fill.fore_color.rgb = RGBColor(220, 252, 231)  # Green bg
                status_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(22, 163, 74)
            elif status == 'IN_PROGRESS':
                status_cell.text = "⏳ Active"
                status_cell.fill.solid()
                status_cell.fill.fore_color.rgb = RGBColor(254, 249, 195)  # Yellow bg
                status_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(161, 98, 7)
            else:
                status_cell.text = "○ Pending"
                status_cell.fill.solid()
                status_cell.fill.fore_color.rgb = RGBColor(243, 244, 246)  # Gray bg
                status_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            
            status_cell.text_frame.paragraphs[0].font.size = Pt(9)
            status_cell.text_frame.paragraphs[0].font.bold = True
            status_cell.vertical_anchor = 1
            
            # Alternating row colors
            if idx % 2 == 0:
                for col in [0, 1, 2, 3]:  # Not status column
                    table.cell(idx, col).fill.solid()
                    table.cell(idx, col).fill.fore_color.rgb = RGBColor(249, 250, 251)
    
    def _add_risks_slide_app_format(self, prs, risks):
        """Add risks slide - clean table format"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "⚠️ Risk Analysis"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
        
        if not risks:
            no_data = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            no_data.text_frame.text = "No active risks"
            no_data.text_frame.paragraphs[0].font.size = Pt(18)
            no_data.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            no_data.text_frame.paragraphs[0].font.color.rgb = RGBColor(107, 114, 128)
            return
        
        # Sort by severity
        severity_order = {'HIGH': 0, 'MEDIUM': 1, 'LOW': 2}
        sorted_risks = sorted(
            risks,
            key=lambda r: severity_order.get(r.get('severity', 'LOW'), 3)
        )[:10]
        
        # Create table
        rows = len(sorted_risks) + 1
        cols = 5  # Severity, Description, Project, Status, Mitigation
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.0),
            Inches(9), Inches(6)
        ).table
        
        # Set column widths
        table.columns[0].width = Inches(1.0)  # Severity
        table.columns[1].width = Inches(2.8)  # Description
        table.columns[2].width = Inches(1.8)  # Project
        table.columns[3].width = Inches(1.0)  # Status
        table.columns[4].width = Inches(2.4)  # Mitigation
        
        # Header row
        headers = ['Severity', 'Risk Description', 'Project', 'Status', 'Mitigation']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(249, 250, 251)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
        
        # Data rows
        for idx, risk in enumerate(sorted_risks, 1):
            # Severity
            sev_cell = table.cell(idx, 0)
            severity = str(risk.get('severity', 'LOW'))
            sev_cell.text = severity
            sev_cell.text_frame.paragraphs[0].font.size = Pt(9)
            sev_cell.text_frame.paragraphs[0].font.bold = True
            sev_cell.vertical_anchor = 1
            
            if severity == 'HIGH':
                sev_cell.fill.solid()
                sev_cell.fill.fore_color.rgb = RGBColor(254, 226, 226)  # Light red
                sev_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(185, 28, 28)
            elif severity == 'MEDIUM':
                sev_cell.fill.solid()
                sev_cell.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Light yellow
                sev_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(161, 98, 7)
            else:
                sev_cell.fill.solid()
                sev_cell.fill.fore_color.rgb = RGBColor(243, 244, 246)  # Light gray
                sev_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            
            # Description
            desc_cell = table.cell(idx, 1)
            desc_cell.text = str(risk.get('description', 'No description'))[:70]
            desc_cell.text_frame.paragraphs[0].font.size = Pt(9)
            desc_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            desc_cell.text_frame.word_wrap = True
            desc_cell.vertical_anchor = 1
            
            # Project
            proj_cell = table.cell(idx, 2)
            proj_cell.text = str(risk.get('project', ''))[:30]
            proj_cell.text_frame.paragraphs[0].font.size = Pt(9)
            proj_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            proj_cell.text_frame.word_wrap = True
            proj_cell.vertical_anchor = 1
            
            # Status
            status_cell = table.cell(idx, 3)
            status_cell.text = str(risk.get('status', 'OPEN'))[:12]
            status_cell.text_frame.paragraphs[0].font.size = Pt(9)
            status_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            status_cell.vertical_anchor = 1
            
            # Mitigation
            mit_cell = table.cell(idx, 4)
            mit_cell.text = str(risk.get('mitigation', 'No mitigation plan'))[:60]
            mit_cell.text_frame.paragraphs[0].font.size = Pt(9)
            mit_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            mit_cell.text_frame.word_wrap = True
            mit_cell.vertical_anchor = 1
            
            # Alternating row colors
            if idx % 2 == 0:
                for col in [1, 2, 3, 4]:  # Not severity column
                    table.cell(idx, col).fill.solid()
                    table.cell(idx, col).fill.fore_color.rgb = RGBColor(249, 250, 251)
    
    def _add_changes_slide_app_format(self, prs, changes):
        """Add change management slide matching app format - table with old→new dates"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "📅 Change Management"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
        
        # Sort by date and take most recent
        sorted_changes = sorted(
            changes,
            key=lambda c: c.get('date', '2025-01-01'),
            reverse=True
        )[:10]
        
        if not sorted_changes:
            return
        
        # Create table matching app format
        rows = len(sorted_changes) + 1
        cols = 5  # Date, Project, Schedule Change (old→new), Reason, Impact
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.0),
            Inches(9), Inches(6)
        ).table
        
        # Set column widths
        table.columns[0].width = Inches(1.2)  # Date
        table.columns[1].width = Inches(1.8)  # Project
        table.columns[2].width = Inches(2.5)  # Schedule Change
        table.columns[3].width = Inches(2.3)  # Reason
        table.columns[4].width = Inches(1.2)  # Impact
        
        # Header row with gray background
        headers = ["Change Date", "Project", "Schedule Change", "Reason", "Impact"]
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(249, 250, 251)  # Gray-50
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Data rows
        for idx, change in enumerate(sorted_changes, 1):
            # Change date
            try:
                change_date = datetime.strptime(str(change.get('date', '')), '%Y-%m-%d')
                date_str = change_date.strftime('%b %d, %Y')
            except:
                date_str = str(change.get('date', ''))
            
            cell = table.cell(idx, 0)
            cell.text = date_str
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            
            # Project
            project = str(change.get('project', ''))[:25]
            cell = table.cell(idx, 1)
            cell.text = project
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            
            # Schedule Change (old → new) with color coding
            try:
                old_date = datetime.strptime(str(change.get('old_date', '')), '%Y-%m-%d')
                old_str = old_date.strftime('%b %d')
            except:
                old_str = str(change.get('old_date', ''))[:10]
            
            try:
                new_date = datetime.strptime(str(change.get('new_date', '')), '%Y-%m-%d')
                new_str = new_date.strftime('%b %d')
            except:
                new_str = str(change.get('new_date', ''))[:10]
            
            cell = table.cell(idx, 2)
            tf = cell.text_frame
            tf.clear()
            
            # Old date (red, strikethrough)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = old_str
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(220, 38, 38)  # Red
            
            # Arrow
            run = p.add_run()
            run.text = " → "
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(107, 114, 128)
            
            # New date (green, bold)
            run = p.add_run()
            run.text = new_str
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(22, 163, 74)  # Green
            run.font.bold = True
            
            # Reason
            reason = str(change.get('reason', ''))[:45]
            cell = table.cell(idx, 3)
            cell.text = reason
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
            cell.text_frame.word_wrap = True
            
            # Impact
            impact = str(change.get('impact', ''))[:20]
            cell = table.cell(idx, 4)
            cell.text = impact
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(107, 114, 128)
            cell.text_frame.word_wrap = True
            
            # Add subtle row background (alternating)
            if idx % 2 == 0:
                for col in range(cols):
                    table.cell(idx, col).fill.solid()
                    table.cell(idx, col).fill.fore_color.rgb = RGBColor(249, 250, 251)

    # ------------------------------------------------------------------
    # Financial Slides (FGSI integration)
    # ------------------------------------------------------------------

    def _add_financial_slides(self, prs):
        """Add financial summary slides if financial data exists."""
        try:
            from repositories.financial_repository import FinancialRepository
            repo = FinancialRepository()
            summary = repo.get_financial_summary()

            if not summary or (summary.get("total_revenue_target", 0) == 0
                               and summary.get("total_cost_budget", 0) == 0):
                logger.info("No financial data — skipping financial slides")
                return

            self._add_financial_kpi_slide(prs, summary)
            self._add_financial_risk_slide(prs, repo)
            logger.info("Financial slides added to PowerPoint export")

        except ImportError:
            logger.debug("Financial repository not available — skipping financial slides")
        except Exception as exc:
            logger.warning(f"Could not generate financial slides: {exc}")

    def _add_financial_kpi_slide(self, prs, summary):
        """Financial KPI summary slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        p.text = "Financial Governance Summary"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 41, 55)

        kpi_data = [
            ("Revenue",
             "\u00a3{:,.0f}".format(summary.get("total_revenue_actual", 0)),
             "Target: \u00a3{:,.0f}".format(summary.get("total_revenue_target", 0)),
             "{:+.1f}%".format(summary.get("revenue_variance_pct", 0))),
            ("Cost",
             "\u00a3{:,.0f}".format(summary.get("total_cost_actual", 0)),
             "Budget: \u00a3{:,.0f}".format(summary.get("total_cost_budget", 0)),
             "{:+.1f}%".format(summary.get("cost_variance_pct", 0))),
            ("Margin",
             "{:.1f}%".format(summary.get("actual_margin", 0) or 0),
             "Target: {:.1f}%".format(summary.get("target_margin", 0) or 0),
             "On Track" if summary.get("on_track", True) else "Off Track"),
            ("Risk Score",
             str(summary.get("financial_risk_score", 0) or 0),
             "{} programs".format(summary.get("program_count", 0)),
             ""),
        ]

        kpi_rows = len(kpi_data) + 1
        kpi_cols = 4
        tbl_shape = slide.shapes.add_table(
            kpi_rows, kpi_cols, Inches(0.5), Inches(1.2),
            Inches(9), Inches(0.5 * kpi_rows))
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(2)
        tbl.columns[1].width = Inches(2.5)
        tbl.columns[2].width = Inches(2.5)
        tbl.columns[3].width = Inches(2)

        for i, h in enumerate(["Metric", "Actual", "Target/Budget", "Variance"]):
            cell = tbl.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(59, 130, 246)

        for ri, (metric, actual, target, variance) in enumerate(kpi_data):
            for ci, val in enumerate([metric, actual, target, variance]):
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(55, 65, 81)
                if ci == 0:
                    cell.text_frame.paragraphs[0].font.bold = True

    def _add_financial_risk_slide(self, prs, repo):
        """Financial risks summary slide."""
        fin_risks = repo.list_financial_risks(status="OPEN")
        if not fin_risks:
            return

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        p.text = "Financial Risks ({} Open)".format(len(fin_risks))
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 41, 55)

        fr_rows = min(len(fin_risks), 10) + 1
        fr_cols = 5
        fr_shape = slide.shapes.add_table(
            fr_rows, fr_cols, Inches(0.5), Inches(1.2),
            Inches(9), Inches(0.5 * fr_rows))
        fr_tbl = fr_shape.table
        fr_tbl.columns[0].width = Inches(1.2)
        fr_tbl.columns[1].width = Inches(3.5)
        fr_tbl.columns[2].width = Inches(1)
        fr_tbl.columns[3].width = Inches(1.5)
        fr_tbl.columns[4].width = Inches(1.8)

        for i, h in enumerate(["Type", "Description", "Severity", "Program", "Mitigation"]):
            cell = fr_tbl.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(220, 38, 38)

        for ri, risk in enumerate(fin_risks[:10]):
            vals = [
                (risk.get("risk_type", "") or "").replace("_", " ").title(),
                (risk.get("description", "") or "")[:80],
                (risk.get("severity", "") or "").upper(),
                risk.get("program_id") or "Portfolio",
                (risk.get("mitigation", "") or "")[:60],
            ]
            for ci, val in enumerate(vals):
                cell = fr_tbl.cell(ri + 1, ci)
                cell.text = val
                cell.text_frame.paragraphs[0].font.size = Pt(8)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(75, 85, 99)
                cell.text_frame.word_wrap = True
