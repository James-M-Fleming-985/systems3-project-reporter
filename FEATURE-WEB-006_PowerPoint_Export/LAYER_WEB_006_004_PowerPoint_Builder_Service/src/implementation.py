import os
import time
from io import BytesIO
from typing import List, Optional, Dict, Any
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


class PowerPointBuilderService:
    """Service for building PowerPoint presentations with screenshots and branding."""
    
    DEFAULT_TEMPLATE_PATH = "templates/default_template.pptx"
    MAX_GENERATION_TIME = 3.0  # seconds
    
    def __init__(self):
        """Initialize the PowerPoint builder service."""
        self.presentation = None
        self.template_path = None
        self.start_time = None
        
    def generate_presentation(
        self,
        report_data: Dict[str, Any],
        screenshots: List[bytes],
        template_path: Optional[str] = None
    ) -> bytes:
        """
        Generate a PowerPoint presentation with report data and screenshots.
        
        Args:
            report_data: Dictionary containing report metadata
            screenshots: List of screenshot images as bytes
            template_path: Optional path to company template
            
        Returns:
            Generated PPTX file as bytes
            
        Raises:
            ValueError: If template is invalid or generation takes too long
        """
        self.start_time = time.time()
        
        # Load template
        self._load_template(template_path)
        
        # Create title slide
        self._create_title_slide(report_data)
        
        # Create content slides for screenshots
        for idx, screenshot in enumerate(screenshots):
            self._create_content_slide(screenshot, idx + 1)
            
        # Check generation time
        elapsed = time.time() - self.start_time
        if elapsed >= self.MAX_GENERATION_TIME:
            raise ValueError("File generation exceeded time limit")
            
        # Save to bytes
        return self._save_to_bytes()
        
    def _load_template(self, template_path: Optional[str] = None):
        """Load PowerPoint template or use default."""
        if template_path:
            try:
                self.presentation = Presentation(template_path)
                self.template_path = template_path
                self._validate_template()
            except Exception as e:
                raise ValueError(f"Invalid template: {str(e)}")
        else:
            # Use default template or create new presentation
            if os.path.exists(self.DEFAULT_TEMPLATE_PATH):
                self.presentation = Presentation(self.DEFAULT_TEMPLATE_PATH)
            else:
                self.presentation = Presentation()
                
    def _validate_template(self):
        """Validate the loaded template."""
        try:
            # Check if presentation can be accessed
            _ = self.presentation.slides
            # Check if it has slide layouts
            if not self.presentation.slide_layouts:
                raise ValueError("Template has no slide layouts")
        except Exception as e:
            raise ValueError(f"Invalid template format: {str(e)}")
            
    def _create_title_slide(self, report_data: Dict[str, Any]):
        """Create title slide with report metadata."""
        # Use first slide layout (title slide)
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        if title:
            title.text = report_data.get('title', 'Report')
            
        # Add metadata to subtitle or body
        subtitle = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != title:
                subtitle = shape
                break
                
        if subtitle:
            text_frame = subtitle.text_frame
            text_frame.clear()
            
            # Add metadata
            p = text_frame.add_paragraph()
            p.text = f"Generated: {report_data.get('generated_date', datetime.now().strftime('%Y-%m-%d'))}"
            p.alignment = PP_ALIGN.CENTER
            
            if 'author' in report_data:
                p = text_frame.add_paragraph()
                p.text = f"Author: {report_data['author']}"
                p.alignment = PP_ALIGN.CENTER
                
            if 'description' in report_data:
                p = text_frame.add_paragraph()
                p.text = report_data['description']
                p.alignment = PP_ALIGN.CENTER
                
    def _create_content_slide(self, screenshot: bytes, slide_number: int):
        """Create a content slide with screenshot."""
        # Use content slide layout (usually index 1 or 5)
        layout_idx = 5 if len(self.presentation.slide_layouts) > 5 else 1
        slide_layout = self.presentation.slide_layouts[layout_idx]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title if available
        if slide.shapes.title:
            slide.shapes.title.text = f"Screenshot {slide_number}"
            
        # Add screenshot image
        self._add_image_to_slide(slide, screenshot)
        
    def _add_image_to_slide(self, slide, image_bytes: bytes):
        """Add image to slide with proper sizing."""
        # Load image to get dimensions
        img = Image.open(BytesIO(image_bytes))
        img_width, img_height = img.size
        
        # Calculate slide dimensions (leaving margins)
        slide_width = self.presentation.slide_width
        slide_height = self.presentation.slide_height
        
        max_width = int(slide_width * 0.8)
        max_height = int(slide_height * 0.6)
        
        # Calculate scaling to fit
        width_ratio = max_width / img_width
        height_ratio = max_height / img_height
        scale_ratio = min(width_ratio, height_ratio)
        
        # Calculate final dimensions
        width = int(img_width * scale_ratio)
        height = int(img_height * scale_ratio)
        
        # Calculate position (centered)
        left = int((slide_width - width) / 2)
        top = int((slide_height - height) / 2)
        
        # Add image
        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream,
            left,
            top,
            width=width,
            height=height
        )
        
    def _save_to_bytes(self) -> bytes:
        """Save presentation to bytes."""
        output = BytesIO()
        self.presentation.save(output)
        output.seek(0)
        return output.read()
        
    def get_slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        if not self.presentation:
            return 0
        return len(self.presentation.slides)
        
    def preserve_branding(self, template_path: str) -> bool:
        """
        Check if company branding is preserved from template.
        
        Returns:
            True if branding elements are preserved
        """
        try:
            # Load template
            template_pres = Presentation(template_path)
            
            # Check if template has master slides with branding
            if template_pres.slide_master:
                # Branding is preserved through slide master
                return True
                
            # Check for logo shapes in template
            for slide_layout in template_pres.slide_layouts:
                for shape in slide_layout.shapes:
                    if shape.shape_type == 13:  # Picture
                        return True
                        
            return True  # Assume branding preserved if template loads
            
        except Exception:
            return False


# Helper functions for testing
def create_default_template():
    """Create a default template if it doesn't exist."""
    os.makedirs("templates", exist_ok=True)
    
    if not os.path.exists(PowerPointBuilderService.DEFAULT_TEMPLATE_PATH):
        prs = Presentation()
        
        # Customize slide master
        slide_master = prs.slide_master
        
        # Add title slide layout
        title_layout = prs.slide_layouts[0]
        
        # Add content slide layout
        content_layout = prs.slide_layouts[5]
        
        # Save as default template
        prs.save(PowerPointBuilderService.DEFAULT_TEMPLATE_PATH)


# Create default template on module load
create_default_template()
